#!/usr/bin/env python3
"""Generate SunMRRC SDD V3.4 architecture overview PPTX — 12 slides, dark minimal."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Theme ────────────────────────────────────────────────────────
BG      = RGBColor(0x0A, 0x0A, 0x15)   # deep navy black
BG_CARD = RGBColor(0x10, 0x15, 0x25)   # card surfaces
ACCENT  = RGBColor(0x5B, 0x9B, 0xD5)   # blue accent
GREEN   = RGBColor(0x6B, 0xCF, 0x7F)   # success green
ORANGE  = RGBColor(0xE8, 0xA8, 0x40)   # warning amber
RED     = RGBColor(0xCF, 0x6A, 0x6A)   # danger red
PURPLE  = RGBColor(0x8A, 0x6A, 0xCF)   # purple
WHITE   = RGBColor(0xE8, 0xEC, 0xF0)   # primary text
GRAY    = RGBColor(0x6A, 0x7A, 0x8A)   # secondary text
DIM     = RGBColor(0x3A, 0x4A, 0x5A)   # dim/border
BG_INNER= RGBColor(0x14, 0x1C, 0x2E)   # nested card

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ───────────────────────────────────────────────────────
def dark_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide

def add_text(slide, left, top, width, height, text, font_size=12,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Inter'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_card(slide, left, top, w, h, fill=BG_CARD, border=DIM, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.5)
    return shape

def add_rich_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_run(tf, text, size=12, color=WHITE, bold=False, name='Inter', new_para=True):
    if new_para and len(tf.paragraphs) > 0 and tf.paragraphs[-1].text != '':
        p = tf.add_paragraph()
    elif len(tf.paragraphs) == 0 or tf.paragraphs[-1].text == '':
        p = tf.paragraphs[0] if len(tf.paragraphs) > 0 else tf.add_paragraph()
    else:
        p = tf.paragraphs[-1]
    if p.text == '':
        p.text = text
        run = p.runs[0] if p.runs else p.add_run()
    else:
        run = p.add_run()
        run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name
    return run

def slide_number(slide, num):
    add_text(slide, 12.2, 7.05, 1, 0.35, str(num), font_size=10, color=DIM, alignment=PP_ALIGN.RIGHT)

def footer(slide, text):
    add_text(slide, 0.5, 7.05, 4, 0.35, text, font_size=8, color=DIM)

def section_title(slide, text, y=0.4):
    add_text(slide, 0.8, y, 10, 0.55, text, font_size=28, color=WHITE, bold=True)

def separator(slide, y):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.8), Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

def tag_label(slide, left, top, text, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(len(text)*0.1+0.3), Inches(0.28))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(color[0]//6, color[1]//6, color[2]//6)
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(7); p.font.color.rgb = color; p.font.bold = True; p.font.name = 'JetBrains Mono'

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
add_text(s, 0.8, 1.8, 11, 1.0, 'SunMRRC', font_size=56, color=WHITE, bold=True)
add_text(s, 0.8, 2.75, 11, 0.6, 'SunSDR2 DX Mobile Remote Radio Control', font_size=22, color=ACCENT)
separator(s, 3.5)
add_text(s, 0.8, 3.85, 11, 0.5, 'Architecture Overview  ·  SDD V3.4  ·  12 Architecture Decisions', font_size=16, color=GRAY)
tf = add_rich_box(s, 0.8, 4.5, 11, 1.0)
add_run(tf, '📡  HTTPS/WSS  ·  🎛️  DSP  ·  🎤  TX Voice (Hilbert SSB + tanh limiter)  ·  📊  0x1F00 Telemetry  ·  🔒  PTT Safety', size=13, color=GRAY)
add_run(tf, '', size=6, color=GRAY)
add_run(tf, 'Python 3.12  ·  FastAPI/Uvicorn  ·  NumPy/SciPy  ·  Web Audio API  ·  vanilla JavaScript  ·  WDSP', size=12, color=DIM)
add_text(s, 0.8, 6.5, 11, 0.4, '2026-06-26  ·  SunMRRC V1.0  ·  https://radio.vlsc.net:8889', font_size=11, color=DIM)
slide_number(s, 1)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — System at a Glance
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'System at a Glance')
separator(s, 1.0)
add_text(s, 0.8, 1.15, 11, 0.4, 'Four-layer architecture: browser → server → DSP → radio hardware', font_size=13, color=GRAY)

# Layer boxes
layers = [
    (1.7, '📱  Mobile Browser', 'iOS Safari / Chrome · Web Audio · WASM Opus · Canvas Waterfall\nTX EQ (4 presets) · AudioWorklet 48k→16k · PTT Safety · Service Worker', ACCENT),
    (3.05, '🖥️  SunMRRC Server  (FastAPI + Uvicorn)', 'Static files · 5 WebSocket endpoints · Auth middleware · TLS (auto-detect)\nIQ processing loop · TX pacer thread · /api/band_power + /api/mem_channels', GREEN),
    (4.65, '🧮  Shared DSP Modules  (web_control/)', 'SunSDR2DXClient · StreamProcessor · TXModulator (Hilbert SSB)\nAudioDemodulator · SpectrumProcessor · Opus Codec · WDSP Wrapper', ORANGE),
    (6.0, '📡  SunSDR2 DX Hardware', 'UDP :50001 Control (DRIVE 0x0017) · UDP :50002 IQ Stream (0xFFFD/0xFFFE)\n0x1F00 Telemetry (forward W · supply V · PA temp °C)  —  192.168.16.200', RED),
]
for y, title, desc, color in layers:
    add_card(s, 0.8, y, 11.7, 1.1, fill=BG_CARD, border=color)
    add_text(s, 1.0, y+0.05, 11.2, 0.3, title, font_size=14, color=color, bold=True)
    add_text(s, 1.0, y+0.35, 11.2, 0.65, desc, font_size=10, color=GRAY)

# Arrows between layers
for y in [2.82, 4.42, 5.78]:
    add_text(s, 6.3, y, 1, 0.2, '↕', font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)

add_text(s, 0.8, 7.15, 11.5, 0.3, 'HTTPS: certs/fullchain.pem + radio.vlsc.net.key  ·  WSS when page protocol = HTTPS  ·  DISABLE_SSL=1 for HTTP dev', font_size=9, color=DIM)
slide_number(s, 2)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — WebSocket Architecture
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'WebSocket Architecture  —  5 Endpoints')
separator(s, 1.0)
add_text(s, 0.8, 1.15, 11, 0.4, 'All connections require auth token (query param) · Tagged binary framing for audio · Bidirectional only where needed', font_size=13, color=GRAY)

endpoints = [
    ('/WSCTRX', 'Bidirectional Text', 'Frequency · Mode · PTT · Tune · Drive · Gain\nFilter · WDSP · Opus toggle · Sample Rate\nPING→PONG latency meter', '⚡ Control', ACCENT),
    ('/WSaudioRX', 'Server → Client Binary', '0x00=PCM Int16 · 0x01=Opus (28kbps VBR)\n16kHz mono · Server-side encode via ctypes\nlibopus fallback → PCM transparent', '🔊 RX Audio', GREEN),
    ('/WSaudioTX', 'Client → Server Binary', '0x00=PCM Int16 · 0x01=Opus · 16kHz mono\nBrowser AudioWorklet 48k→16k → WASM Opus\n→ TXModulator Hilbert SSB → 0xFFFD IQ', '🎤 TX Mic', ORANGE),
    ('/WSspectrum', 'Server → Client Binary', '512 × uint8 bytes/frame (~38 Hz)\n0=-120dB · 255=0dB · Client accumulates\n10 frames → adaptive noise floor → waterfall', '📊 Spectrum', PURPLE),
    ('/WSATR1000', 'Bidirectional JSON', 'Antenna tuner proxy (placeholder)\nFrontend hooks exist · Backend stub only\nAccepts connections · No real HW interface', '📡 Tuner', DIM),
]
y = 1.7
for ep, direction, detail, icon_label, color in endpoints:
    add_card(s, 0.8, y, 11.7, 0.95, fill=BG_CARD, border=color)
    tag_label(s, 1.0, y+0.08, icon_label, color)
    add_text(s, 2.6, y+0.05, 3, 0.25, ep, font_size=13, color=color, bold=True)
    add_text(s, 2.6, y+0.28, 3, 0.2, direction, font_size=9, color=DIM)
    add_text(s, 5.8, y+0.05, 6.5, 0.85, detail, font_size=10, color=GRAY)
    y += 1.05

add_text(s, 0.8, 7.0, 11.5, 0.3, 'Client: controls.js + mobile.js + tx_button.js + modules/  ·  Server: ws_ctrl() / ws_audio_rx() / ws_audio_tx() / ws_spectrum()', font_size=9, color=DIM)
slide_number(s, 3)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — RX Signal Chain
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'RX Signal Chain  —  IQ to Audio')
separator(s, 1.0)
add_text(s, 0.8, 1.15, 11, 0.4, 'RX audio path: device IQ stream → DSP demodulation → Opus/PCM encode → browser playback', font_size=13, color=GRAY)

stages_rx = [
    ('① UDP Intake', '0xFFFE packets\n78,125 Hz IQ\n24-bit signed LE', ACCENT),
    ('② Decode', '200 samples/pkt\nnormalize to\ncomplex64', ACCENT),
    ('③ Demodulate', 'IF shift +30.5kHz\nUSB/LSB bandpass\n→ real audio', GREEN),
    ('④ Decimate', '5:1 → 15,625 Hz\nAGC (WDSP or\nbuilt-in fallback)', GREEN),
    ('⑤ WDSP (opt.)', 'NR2 · NB · ANF\nAGC SLOW/MED/\nFAST · Notches', PURPLE),
    ('⑥ Encode', 'Resample→16kHz\nOpus 28kbps VBR\nor Int16 PCM', ORANGE),
    ('⑦ Broadcast', '/WSaudioRX\ntagged frames\n→ browser', ORANGE),
    ('⑧ Play', 'AudioContext\nWASM Opus decode\nor PCM→Float32', ACCENT),
]
x = 0.5
for title, desc, color in stages_rx:
    add_card(s, x, 1.7, 1.45, 1.6, fill=BG_INNER, border=color)
    add_text(s, x+0.08, 1.78, 1.3, 0.3, title, font_size=10, color=color, bold=True)
    add_text(s, x+0.08, 2.1, 1.3, 1.1, desc, font_size=9, color=GRAY)
    if x > 0.5:
        add_text(s, x-0.3, 2.3, 0.3, 0.3, '→', font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)
    x += 1.55

# RX details
tf = add_rich_box(s, 0.8, 3.6, 11.5, 3.2)
add_run(tf, 'Modulation Path', size=14, color=ACCENT, bold=True)
add_run(tf, 'IQ → IF shift (30.5 kHz) → complex SSB bandpass (FIR, hamming window, 60dB attn) → real extraction → 5:1 decimate → built-in AGC', size=11, color=GRAY)
add_run(tf, '', size=6, color=GRAY)
add_run(tf, 'Waterfall Pipeline', size=14, color=ACCENT, bold=True)
add_run(tf, 'IQ → FFT (2048-point, hanning) → dB clip [-120, 0] → uint8 quantize → /WSspectrum (512B/frame, ~38Hz) → client accumulates 10 frames → 30th percentile noise floor → blue bias + contrast gain → black→blue→cyan→yellow→red colour ramp', size=11, color=GRAY)
add_run(tf, '', size=6, color=GRAY)
add_run(tf, 'S-Meter Pipeline', size=14, color=ACCENT, bold=True)
add_run(tf, 'FFT → 90th percentile → S9=9+(p90+73)/6 → asymmetric EMA (attack α=0.5, release α=0.15) → stable needle, fast-rise slow-decay', size=11, color=GRAY)
add_run(tf, '', size=6, color=GRAY)
add_run(tf, 'Audio Codec', size=14, color=ACCENT, bold=True)
add_run(tf, 'Server: direct ctypes libopus — bitrate via opus_encode() max_data_bytes cap (avoids arm64 variadic ctl bug). Client: WASM OpusDecoder. Each /WSaudioRX frame prefixed with 1-byte codec tag (0x00=PCM, 0x01=Opus). Default Opus; switchable via Audio Codec menu. Falls back to PCM if libopus missing.', size=11, color=GRAY)
slide_number(s, 4)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — TX Voice Chain
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'TX Voice Chain  —  Mic to RF')
separator(s, 1.0)
add_text(s, 0.8, 1.15, 11, 0.4, 'End-to-end TX path: browser Web Audio pipeline → Opus → WebSocket → Hilbert SSB → device IQ stream', font_size=13, color=GRAY)

# Client side
add_card(s, 0.5, 1.7, 12.3, 2.5, fill=BG_CARD, border=ACCENT)
add_text(s, 0.7, 1.75, 6, 0.3, '📱  Client: Web Audio TX EQ Pipeline  (48 kHz)', font_size=14, color=ACCENT, bold=True)
client_stages = [
    ('Mic\nSource', '48kHz\nmono'),
    ('Preamp\n×1.5', '+3.5dB\nheadroom'),
    ('AntiAlias\n×2', 'LP 4.5kHz\nQ=0.707'),
    ('EQ Low\n350Hz pk', '+6~12dB\nband'),
    ('EQ Mid\n1500Hz pk', '+8~12dB\nclarity'),
    ('EQ High\n2700Hz sh', '-6~-18dB\nhiss cut'),
    ('Comp\n3:1', 'thr=-24\nknee=12'),
    ('Makeup\n×1.6', '+4dB\npost-comp'),
    ('Noise\nGate', 'thr=-50dB\nslow rel'),
    ('Audio\nWorklet', '48k→16k\nInt16'),
    ('Opus\nEncode', '28kbps\nVBR FEC'),
    ('WSaudioTX\n0x01 tag', '20ms\nframes'),
]
cx = 0.55
for title, desc in client_stages:
    add_card(s, cx, 2.15, 0.92, 0.95, fill=BG_INNER, border=ACCENT)
    add_text(s, cx+0.03, 2.2, 0.86, 0.4, title, font_size=7.5, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, cx+0.03, 2.6, 0.86, 0.45, desc, font_size=6.5, color=GRAY, alignment=PP_ALIGN.CENTER)
    if cx > 0.55:
        add_text(s, cx-0.15, 2.45, 0.2, 0.2, '→', font_size=10, color=DIM, alignment=PP_ALIGN.CENTER)
    cx += 1.0

# EQ presets
add_text(s, 0.7, 3.18, 11.5, 0.2, 'EQ Presets:  DEFAULT (+6/+8/-6dB) · MEDIUM (+9/+10/-12dB) · STRONG (+12/+12/-18dB)  |  RAGCHEW: lowCut 150Hz · midCut 500Hz -2dB · presence 2.4kHz +3dB · highCut 3kHz · 3:1 comp', font_size=8, color=DIM)

# Server side
add_card(s, 0.5, 4.35, 12.3, 1.8, fill=BG_CARD, border=GREEN)
add_text(s, 0.7, 4.4, 6, 0.3, '🖥️  Server: TXModulator DSP Pipeline  (dsp.py)', font_size=14, color=GREEN, bold=True)
srv_stages = [
    ('Resampler', '16k→15625Hz\nfractional\ncontinuous', GREEN),
    ('Hilbert SSB\noverlap-save', 'MARGIN=256\n80-sample hops\nphase-clean', GREEN),
    ('Upsample\n×2.5', '→39063Hz\n200 IQ/pkt\nlinear interp', GREEN),
    ('Drive Gain\n×3.0×drive%', 'fixed makeup\nNOT a power\ncontrol', ORANGE),
    ('tanh(1.0)\nsoft limiter', 'smooth ceiling\n~4% engagement\nsafety net', RED),
    ('24-bit IQ\npacking', 'vectorized\nnumpy ops\nGIL-free', GREEN),
    ('Jitter\nBuffer', 'prime 24pkts\nreprime 12\nhysteresis', GREEN),
    ('TX Pacer\nThread', '5.12ms/pkt\nadaptive ±15%\n0xFFFD', ACCENT),
]
sx = 0.55
for title, desc, color in srv_stages:
    add_card(s, sx, 4.8, 1.42, 1.2, fill=BG_INNER, border=color)
    add_text(s, sx+0.05, 4.85, 1.32, 0.45, title, font_size=8.5, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, sx+0.05, 5.3, 1.32, 0.65, desc, font_size=7, color=GRAY, alignment=PP_ALIGN.CENTER)
    if sx > 0.55:
        add_text(s, sx-0.15, 5.2, 0.2, 0.2, '→', font_size=10, color=DIM, alignment=PP_ALIGN.CENTER)
    sx += 1.5

# Device
add_text(s, 0.8, 6.3, 11.5, 0.2, '📡  Device:  0xFFFD UDP packets → 195 pkt/s (39063 Hz IQ) → DRIVE command (0x0017) per-band power → PA (no firmware ALC) → Antenna', font_size=10, color=RED)
add_text(s, 0.8, 6.55, 11.5, 0.2, '⚡  Tune mode bypasses this entire chain:  pre-computed IQ with TX_TUNE_SCALE=0.35  →  ~10W continuous (safe for PA thermal)  →  no client/server audio path involved', font_size=10, color=ORANGE)
add_text(s, 0.8, 6.85, 11.5, 0.25, '🔑  Power is set ONLY at the DEVICE (0x0017).  Client gain = healthy digital level.  Server gain = lifts into tanh knee.  Turn drive ↑ for more RF power.', font_size=10, color=WHITE, bold=True)
slide_number(s, 5)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — TX Gain Staging (AD-012)
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'TX Gain Staging  —  AD-012')
separator(s, 1.0)
add_text(s, 0.8, 1.15, 11, 0.5, 'Client preamp was reduced from ×3.0 → ×1.5 on 2026-06-25. Old value saturated the tanh (75% reduction → heavy distortion). New value gives ~4% tanh engagement → clean SSB.', font_size=13, color=GRAY)

# Healthy path
add_card(s, 0.5, 1.8, 12.3, 2.2, fill=BG_CARD, border=GREEN)
add_text(s, 0.7, 1.85, 6, 0.3, '✅  Healthy Gain Staging  (preamp ×1.5)', font_size=14, color=GREEN, bold=True)

stages_healthy = [
    ('Mic Input\n+ Preamp ×1.5\n+ EQ STRONG', 'peak 0.50\nRMS ~0.05', ACCENT),
    ('After\nHilbert SSB\n+~30% peak', 'peak 0.65\nRMS ~0.07', GREEN),
    ('After Drive\nGain ×3.0\n@ 100% drive', 'peak 1.95\nRMS ~0.21', ORANGE),
    ('After\ntanh(1.0)\nsoft limiter', 'peak 0.96\nRMS ~0.19', GREEN),
    ('24-bit IQ\n→ 0xFFFD\n→ SunSDR PA', 'DRIVE 0x0017\nRF power', ACCENT),
]
hx = 0.6
for title, desc, color in stages_healthy:
    add_card(s, hx, 2.25, 2.2, 1.4, fill=BG_INNER, border=color)
    add_text(s, hx+0.1, 2.3, 2, 0.55, title, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, hx+0.1, 2.85, 2, 0.7, desc, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
    if hx > 0.6:
        add_text(s, hx-0.25, 2.75, 0.3, 0.3, '→', font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)
    hx += 2.4

add_text(s, 1.0, 3.7, 10, 0.2, 'tanh barely engages (~4% peak reduction)  —  clean SSB envelope  —  device drive controls actual RF power', font_size=10, color=GREEN, bold=True)

# Distorted path (old)
add_card(s, 0.5, 4.2, 12.3, 1.6, fill=BG_CARD, border=RED)
add_text(s, 0.7, 4.25, 6, 0.3, '❌  Old Distorted Path  (preamp ×3.0)  —  corrected 2026-06-25', font_size=14, color=RED, bold=True)
dx = 0.6
for title, desc, color in [
    ('Mic + Preamp ×3.0\n+ STRONG EQ', 'peak 1.00\n→ AudioWorklet clips', RED),
    ('After Hilbert\n~+30%', 'peak 1.32', RED),
    ('After Drive\nGain ×3.0', 'peak 3.95\n→ 4× over ceiling', RED),
    ('After tanh(1.0)\n→ 75% SQUASHED', 'peak 0.999\nHEAVY DISTORTION', RED),
    ('Far end hears:\n"splattery" noise\non every syllable', '−6dB fix:\npreamp ×1.5', RED),
]:
    add_card(s, dx, 4.55, 2.2, 1.1, fill=BG_INNER, border=color)
    add_text(s, dx+0.1, 4.6, 2, 0.5, title, font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, dx+0.1, 5.1, 2, 0.5, desc, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)
    if dx > 0.6:
        add_text(s, dx-0.25, 4.9, 0.3, 0.3, '→', font_size=14, color=RED, alignment=PP_ALIGN.CENTER)
    dx += 2.4

# Level probe diagnostics
add_card(s, 0.5, 6.0, 12.3, 1.0, fill=BG_INNER, border=DIM)
add_text(s, 0.7, 6.05, 11.5, 0.25, '📊  Level Probe Diagnostics  (1 Hz in server.log during TX)', font_size=12, color=ACCENT, bold=True)
add_text(s, 0.7, 6.3, 11.5, 0.6, 'TX chain in : rms=… peak=…  n=…     ←  Post-decode PCM (before Hilbert, scaled by 1/gain)\nTX chain an : rms=… peak=…  n=…     ←  After Hilbert analytic signal\nTX chain drv: rms=… peak=…  n=…     ←  After drive gain (pre-tanh, full scale)\nTX chain lim: rms=… peak=…  n=…     ←  After tanh(1.0), final IQ envelope  ·  Healthy: in peak ~0.5, drv peak ~2.0, lim peak ~0.96', font_size=9, color=GRAY)
slide_number(s, 6)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Device Telemetry (0x1F00)
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'Device Telemetry  —  0x1F00  (AD-011)')
separator(s, 1.0)
add_text(s, 0.8, 1.15, 11, 0.4, 'Verified 2026-06-25 against a real 40m drive sweep (expert_40m_drive.pcap) × external wattmeter + ExpertSDR3', font_size=13, color=GRAY)

# Packet layout
add_card(s, 0.5, 1.7, 12.3, 0.8, fill=BG_CARD, border=ACCENT)
add_text(s, 0.7, 1.75, 11.5, 0.25, '34-Byte Packet:  0xFF32 header (10B)  +  seq/flags  +  off16(u16)  +  off18(f32)  +  off22(f32)  +  off26(f32)  +  off30(f32)', font_size=12, color=ACCENT, bold=True)
add_text(s, 0.7, 2.1, 11.5, 0.25, 'Sent continuously in RX and TX modes  —  Device simply does not expose reverse power  →  cannot compute SWR  →  use external ATR-1000 for SWR', font_size=10, color=DIM)

# Verified fields (3 columns)
fields = [
    ('✅  off30: f32 Forward Watts', GREEN, 'PEP envelope · Direct float\nMonotonic with DRIVE\n28%→3W  71%→54W\n88%→83W  100%→101W\nMatches ExpertSDR3 ~95W\nNO cubic/linear fit needed'),
    ('✅  off16: u16 Supply Volts ×10', GREEN, '~136 (13.6V) at idle\nSags as power rises:\n0W→13.6  30-50W→13.1\n80-110W→12.9V\nCorrelation with W = -0.79\nTextbook PSU sag curve\nNOT SWR — stable when unkeyed'),
    ('✅  off18: f32 PA Temp °C', GREEN, '~42°C · Barely moves\nMonitors PA thermal\nstate during extended TX\n\nFrontend: getTXTelem:\n  watts, volts, temp, W_int\nVOLT field replaces SWR'),
]
fx = 0.6
for title, color, desc in fields:
    add_card(s, fx, 2.7, 3.9, 2.3, fill=BG_CARD, border=color)
    add_text(s, fx+0.15, 2.78, 3.6, 0.3, title, font_size=14, color=color, bold=True)
    add_text(s, fx+0.15, 3.15, 3.6, 1.8, desc, font_size=10, color=GRAY)
    fx += 4.1

# What was wrong
add_card(s, 0.5, 5.2, 12.3, 1.8, fill=BG_CARD, border=RED)
add_text(s, 0.7, 5.25, 6, 0.3, '⚠  Previously Wrong  (corrected 2026-06-25)', font_size=14, color=ORANGE, bold=True)
wrong_items = [
    ('OLD: off14 u16 + cubic fit\n(pwr_raw-9)³ × 1.91e-5', 'non-monotonic noise\n28%→72, 45%→51, 100%→81\nnever tracked real power'),
    ('OLD: off16 u16/100 as "SWR"\nreading 1.32-1.37', 'actually supply voltage\nstable when NOT keyed → rules out SWR\nPSU sag confirms voltage identity'),
    ('OLD: off22 f32 as "reverse power"', 'actually AVERAGE forward power\nratio to off30 ≈ 3:1 SSB crest factor'),
    ('OLD: off26 f32 as "SWR alternative"', 'always exactly 1.0000\ndevice placeholder · 323+ packets\nRX=TX=TUNE all identical'),
]
wx = 0.6
for title, desc in wrong_items:
    add_card(s, wx, 5.65, 2.9, 1.2, fill=BG_INNER, border=DIM)
    add_text(s, wx+0.1, 5.7, 2.7, 0.65, title, font_size=9, color=RED, bold=True)
    add_text(s, wx+0.1, 6.1, 2.7, 0.65, desc, font_size=8, color=GRAY)
    wx += 3.1

slide_number(s, 7)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Architecture Decisions (12 ADs)
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'Architecture Decisions  —  12 ADs')
separator(s, 1.0)

ads = [
    ('AD-001', 'FastAPI/Uvicorn', 'Async server, native WS, small surface'),
    ('AD-002', 'Direct UDP Control', 'Verified SunSDR protocol, no Hamlib'),
    ('AD-003', 'DSP-Owned Mode', 'Hardware mode-agnostic, SW demod selects'),
    ('AD-004', 'Tagged Dual-Codec Audio', 'Opus+PCM RX and TX, 1-byte tag, no race'),
    ('AD-005', 'Quantized Spectrum', '512 uint8 bytes/frame, compact waterfall'),
    ('AD-006', 'HTTPS/WSS Default', 'iOS secure context, mic+AudioContext require'),
    ('AD-007', 'PTT Release Safety', 'ACK retry, backup s:, watchdog, forced-RX'),
    ('AD-008', 'Optional WDSP', 'Falls back cleanly, UI queries availability'),
    ('AD-009', 'Gap Tracking', 'Document hooks until backend exists or removed'),
    ('AD-010', 'Device DRIVE Power', '0x0017 per-band, sqrt taper byte, re-send QSY'),
    ('AD-011', '0x1F00 Telemetry', 'off30=W, off16=V, off18=°C — no SWR from device'),
    ('AD-012', 'TX Gain Staging', 'Preamp×1.5, drive gain×3.0, tanh@1.0, clean SSB'),
]
y = 1.3
for id, topic, desc in ads:
    tag_color = GREEN if id != 'AD-009' else ORANGE
    add_card(s, 0.5, y, 12.3, 0.42, fill=BG_CARD, border=DIM)
    tag_label(s, 0.65, y+0.07, id, tag_color)
    add_text(s, 1.8, y+0.06, 4, 0.3, topic, font_size=12, color=WHITE, bold=True)
    add_text(s, 6.0, y+0.06, 6.5, 0.3, desc, font_size=10, color=GRAY)
    y += 0.47

slide_number(s, 8)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Key Design Insights
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'Key Design Insights')
separator(s, 1.0)

insights = [
    ('🔑', 'RF Power Control', 'Only device DRIVE (0x0017) sets TX power. Client preamp and server drive gain are fixed — they set a healthy digital level, not RF power. Turn drive ↑ for more power. Never chase power by raising client gain; that saturates the tanh and distorts.', ACCENT),
    ('🎯', 'Gain Staging', 'Mouth → preamp×1.5 → EQ → AudioWorklet → Opus → Server Hilbert → drive gain×3.0 → tanh(1.0) → IQ. Each stage has headroom. tanh barely touches the signal (~4%) — it is a safety net for transients, not the primary limiter. Device drive is the sole power lever.', GREEN),
    ('📡', 'Device Telemetry', '0x1F00 off30 f32 = forward watts (direct float). The device sends NO reverse power. Cannot compute SWR from telemetry. Old "SWR" field (off16 u16/100) was actually the 13.6V supply rail. Use an external tuner for SWR. This was verified against a complete 40m drive sweep.', ORANGE),
    ('⚡', 'Tune is Separate', 'Tune mode bypasses the entire TX voice chain. It uses a pre-computed continuous analytic IQ signal with TX_TUNE_SCALE=0.35 (~10W safe carrier). No client audio, no server modulator gain, no tanh. This is why Tune was always clean even when voice was distorted.', PURPLE),
    ('🛡️', 'PTT Safety', 'Release is more safety-critical than keying. Multiple redundant paths: frontend ACK retry (3×), backup s: command on TX socket, 1s watchdog timeout, backend forced-RX handler. Touch-move auto-release on mobile. A lost setPTT:false must never leave the radio stuck in TX.', RED),
    ('📊', 'Observability', '1 Hz level probes in server.log (TX chain in:/an:/drv:/lim:). Per-packet pacer timing in /tmp/tx_probe.csv. Mic arrival jitter in /tmp/tx_rx_probe.csv. opus_encoder failures logged. Device SUB-ID census logged once. Underrun counter surfaced. Debuggable by design.', ACCENT),
]
y = 1.3
for icon, title, desc, color in insights:
    add_card(s, 0.5, y, 12.3, 0.9, fill=BG_CARD, border=color)
    add_text(s, 0.7, y+0.05, 0.5, 0.35, icon, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, 1.3, y+0.05, 4, 0.3, title, font_size=14, color=color, bold=True)
    add_text(s, 1.3, y+0.35, 11.3, 0.5, desc, font_size=10, color=GRAY)
    y += 1.0

slide_number(s, 9)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Component & Service Model
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'Component & Service Model')
separator(s, 1.0)

# Left: Components
add_card(s, 0.5, 1.3, 6.0, 5.6, fill=BG_CARD, border=ACCENT)
add_text(s, 0.7, 1.35, 5.5, 0.3, 'Key Components', font_size=16, color=ACCENT, bold=True)
comps = [
    ('Backend', [
        'FastAPIApp — routes, startup, static serving',
        'RadioClient (sunsdr_direct.py) — UDP protocol',
        'IQLoop — UDP socket, packet validation, DSP feed',
        'TXModulator — Hilbert SSB, drive gain, tanh, IQ pack',
        'TXPacer — adaptive ±15% pacing, jitter buffer',
        'RxOpusEncoder / TxOpusDecoder — ctypes libopus',
        'BandPowerAPI / MemChannelAPI — JSON persistence',
    ]),
    ('Frontend', [
        'controls.js — WS, codec, waterfall, S-meter',
        'mobile.js — UX, menus, DSP panel, settings',
        'tx_audio_eq.js — 4 EQ presets, compressor, preamp',
        'tx_capture_worklet.js — AudioWorklet 48k→16k',
        'ptt_manager.js — ACK retry, watchdog, state sync',
        'tx_button.js — touch PTT, lock, warm-up frames',
        'opus_codec.js — WASM Opus encode/decode (28kbps)',
    ]),
]
cy = 1.75
for section, items in comps:
    add_text(s, 0.8, cy, 5.3, 0.25, section, font_size=11, color=ACCENT, bold=True)
    cy += 0.28
    for item in items:
        add_text(s, 1.1, cy, 5.2, 0.22, item, font_size=9, color=GRAY)
        cy += 0.24
    cy += 0.15

# Right: Services
add_card(s, 6.8, 1.3, 6.0, 5.6, fill=BG_CARD, border=GREEN)
add_text(s, 7.0, 1.35, 5.5, 0.3, 'Service Portfolio', font_size=16, color=GREEN, bold=True)
services = [
    ('StaticUIService', 'serve mobile PWA assets from static/'),
    ('ControlService', '/WSCTRX command dispatch, state sync'),
    ('RXAudioService', 'IQ→PCM→Opus /WSaudioRX fan-out'),
    ('TXAudioIngressService', '/WSaudioTX mic→TXModulator'),
    ('SpectrumService', 'FFT quantize /WSspectrum fan-out'),
    ('SunSDRDeviceService', 'UDP protocol, boot, heartbeat'),
    ('IQProcessingService', 'receive, decode, feed DSP'),
    ('DSPService', 'demodulation + optional WDSP chain'),
    ('TXPowerService', '0x0017 DRIVE per-band + /api/band_power'),
    ('BandPowerService', 'per-band drive % JSON persistence'),
    ('MemoryChannelService', '/api/mem_channels JSON persistence'),
    ('TLSStartupService', 'cert/key auto-detect, HTTP fallback'),
]
sy = 1.75
for name, desc in services:
    add_text(s, 7.1, sy, 5.5, 0.22, name, font_size=9.5, color=GREEN, bold=True)
    add_text(s, 8.8, sy, 3.8, 0.22, desc, font_size=8.5, color=GRAY)
    sy += 0.28

slide_number(s, 10)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Operational Model
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'Operational Model')
separator(s, 1.0)

add_card(s, 0.5, 1.3, 6.0, 2.8, fill=BG_CARD, border=ACCENT)
add_text(s, 0.7, 1.35, 5.5, 0.3, '🚀  Startup & Restart', font_size=16, color=ACCENT, bold=True)
add_text(s, 0.7, 1.75, 5.5, 2.2, '''restart.sh — background production start
  → Kill old process by cwd (safe)
  → Clear listening port
  → Activate ../venv if present
  → uvicorn with TLS auto-detect
  → Log to server.log

start.sh — simple foreground launcher
restart.sh -f — foreground with live logs
DISABLE_SSL=1 — HTTP debug (no iOS mic)

Env Config:
  DEVICE_HOST=192.168.16.200
  WEB_PORT=8080 (default 8889)
  WEB_PASSWORD=sunmrrc (override with env)
  DISABLE_SSL=1 (force HTTP)''', font_size=10, color=GRAY)

add_card(s, 6.8, 1.3, 6.0, 2.8, fill=BG_CARD, border=GREEN)
add_text(s, 7.0, 1.35, 5.5, 0.3, '🔐  Auth & TLS', font_size=16, color=GREEN, bold=True)
add_text(s, 7.0, 1.75, 5.5, 2.2, '''Password-based session auth:
  → /login → POST /api/auth/login
  → sets sunmrrc_auth cookie (30-day)
  → token passed as ?token= on all WS
  → _auth_tokens server-side set
  → All routes + WS require token

TLS auto-detection:
  → certs/fullchain.pem
  → certs/radio.vlsc.net.key
  → Both present → HTTPS/WSS
  → Missing or DISABLE_SSL=1 → HTTP
  → check_ssl_expiry.sh for monitoring

Connection Matrix:
  Browser ↔ SunMRRC : HTTPS/WSS
  SunMRRC ↔ SunSDR2 DX : UDP :50001/50002
  IPs: 192.168.16.100 ← → 192.168.16.200''', font_size=10, color=GRAY)

add_card(s, 0.5, 4.35, 12.3, 2.5, fill=BG_CARD, border=ORANGE)
add_text(s, 0.7, 4.4, 6, 0.3, '🔍  Verification & Diagnostics', font_size=16, color=ORANGE, bold=True)
diags = [
    ('Verify HTTPS', 'server.log → "sunmrrc https://..."'),
    ('Verify Radio', 'server.log → "SunSDR2DX: True"'),
    ('Verify RX', 'UI shows /WSaudioRX connected, waterfall flowing'),
    ('Verify TX Power', 'W= in server.log from 0x1F00 off30 f32 forward watts'),
    ('Verify TX Quality', 'TX chain lim: peak ~0.96 (not 0.999 → distorted)'),
    ('Verify Gain Staging', 'in peak ~0.5, drv peak ~2.0, lim peak ~0.96 @ 1 Hz'),
    ('Verify PTT Safety', 'getPTT:false after release, underruns near 0'),
    ('Diagnose Jitter', '/tmp/tx_probe.csv (pacer) + /tmp/tx_rx_probe.csv (mic arrival)'),
    ('Diagnose IQ Power', 'TX IQ pwr: peak=… rms=… in server.log at 1 Hz'),
    ('Level Probes', 'TX chain in:/an:/drv:/lim: lines at 1 Hz during TX'),
]
dx, dy = 0.7, 4.85
for title, desc in diags:
    add_text(s, dx, dy, 3.8, 0.22, title, font_size=10, color=ORANGE, bold=True)
    add_text(s, dx+3.9, dy, 2.5, 0.22, desc, font_size=9, color=GRAY)
    dy += 0.23
    if dy > 6.5:
        dy = 4.85; dx = 6.5

slide_number(s, 11)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Known Gaps & Roadmap
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
section_title(s, 'Known Gaps & Future Work')
separator(s, 1.0)

add_card(s, 0.5, 1.3, 6.0, 5.5, fill=BG_CARD, border=RED)
add_text(s, 0.7, 1.35, 5.5, 0.3, '⚠  Known Gaps  (AD-009)', font_size=16, color=RED, bold=True)
gaps = [
    ('/WSATR1000', 'ATR-1000 antenna tuner backend stub only. Frontend hooks and status placeholders exist but no real HW interface. Only available SWR source in the system.'),
    ('Fixed LAN IPs', '192.168.16.100 (local) and 192.168.16.200 (device) hardcoded. Deployment tied to current LAN topology. Should move to env/config.'),
    ('No Device SWR', '0x1F00 has no reverse-power field. Cannot compute SWR from telemetry. External tuner (ATR-1000) is the only SWR source.'),
    ('start.sh URL', 'start.sh still prints HTTP URL even when TLS may be used. Align script message with server TLS behavior.'),
]
gy = 1.75
for title, desc in gaps:
    add_text(s, 0.8, gy, 5.5, 0.25, title, font_size=12, color=RED, bold=True)
    add_text(s, 0.8, gy+0.28, 5.5, 0.65, desc, font_size=9.5, color=GRAY)
    gy += 1.1

add_card(s, 6.8, 1.3, 6.0, 5.5, fill=BG_CARD, border=GREEN)
add_text(s, 7.0, 1.35, 5.5, 0.3, '✅  What Is Complete', font_size=16, color=GREEN, bold=True)
done_items = [
    '📱  Mobile-first UI (PWA, safe areas, manifest)',
    '📡  SunSDR2 DX direct UDP protocol (boot/heartbeat)',
    '🔊  RX audio: tagged dual-codec (Opus/PCM, >10× BW cut)',
    '📊  Real-time waterfall: adaptive noise floor + colour ramp',
    '🎤  TX voice: Hilbert SSB → tanh limiter → 24-bit IQ',
    '⚡  TX power: device DRIVE 0x0017 per-band (Band Power UI)',
    '🎛️  TX EQ: 4 presets, compressor, preamp, anti-alias',
    '🛡️  PTT safety: ACK retry, watchdog, backup s:, forced-RX',
    '🎵  WDSP: NR2, NB, ANF, NF, AGC (4 modes), notches',
    '📈  S-meter: asymmetric EMA smoothing (fast attack/slow decay)',
    '💾  Memory channels + Band power: JSON persistence + API',
    '🔤  Sample rate selector: 39/78/156/312 kHz via 0x0001',
    '🔐  Password auth: cookie + WS token, 30-day validity',
    '📋  Recordings: server-side RX MP3 via ffmpeg pipe',
    '📐  12 ADs documented in SDD V3.4 + 3 SVG diagrams',
]
dy = 1.75
for item in done_items:
    add_text(s, 7.1, dy, 5.5, 0.22, item, font_size=9.5, color=GRAY)
    dy += 0.29

add_text(s, 0.8, 7.1, 11.5, 0.3, 'SunMRRC V1.0 is production-ready for mobile RX/control + TX voice + per-band power. Open gaps: ATR backend, configurable LAN IPs, SWR from external tuner only.', font_size=10, color=DIM)
slide_number(s, 12)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — Thank You / Reference
# ═══════════════════════════════════════════════════════════════════
s = dark_slide()
add_text(s, 0.8, 2.2, 11, 1.0, 'SunMRRC', font_size=56, color=WHITE, bold=True)
add_text(s, 0.8, 3.15, 11, 0.5, 'Architecture Overview  ·  SDD V3.4  ·  June 2026', font_size=20, color=ACCENT)
separator(s, 3.8)

ref_links = [
    ('🌐', 'Live Radio', 'https://radio.vlsc.net:8889'),
    ('📖', 'SDD Docs', 'https://www.vlsc.net/sunmrrc/sdd/'),
    ('🔧', 'GitHub', 'https://github.com/cheenle/sunsdr'),
    ('📐', 'SVG Diagrams', 'https://www.vlsc.net/sunmrrc/sdd/diagrams/'),
]
rx = 0.8
for icon, label, url in ref_links:
    add_text(s, rx, 4.2, 3, 0.3, f'{icon}  {label}', font_size=16, color=WHITE, bold=True)
    add_text(s, rx, 4.55, 3, 0.3, url, font_size=11, color=ACCENT)
    rx += 3.0

# Stats
add_card(s, 0.5, 5.2, 12.3, 1.2, fill=BG_CARD, border=DIM)
stats = [
    ('12', 'Architecture\nDecisions'),
    ('15', 'SDD\nChapters'),
    ('5', 'WebSocket\nEndpoints'),
    ('4', 'TX EQ\nPresets'),
    ('2', 'Audio\nCodecs'),
    ('195', 'TX IQ\npkts/sec'),
    ('~512B', 'Spectrum\nper frame'),
    ('~20kbps', 'Opus\nbandwidth'),
]
sx = 0.8
for num, label in stats:
    add_text(s, sx, 5.35, 1.2, 0.4, num, font_size=24, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, sx, 5.85, 1.2, 0.4, label, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)
    sx += 1.5

add_text(s, 0.8, 6.75, 11, 0.4, 'Python 3.12  ·  FastAPI + Uvicorn  ·  NumPy + SciPy  ·  Web Audio API  ·  vanilla JavaScript  ·  libopus  ·  libwdsp  ·  IBM TeamSD v2.3.2', font_size=10, color=DIM)
add_text(s, 0.8, 7.1, 11, 0.3, 'Document ID: SDD-SUNMRRC-2026-001  ·  2026-06-26  ·  GPL-3.0', font_size=9, color=DIM)
slide_number(s, 13)

# ── Save ──────────────────────────────────────────────────────────
out_path = '/Users/cheenle/HAM/sunsdr/SDD/SunMRRC-Architecture-SDD-V3.4.pptx'
prs.save(out_path)
print(f'✅  Saved: {out_path}')
print(f'Slides: {len(prs.slides)}')
