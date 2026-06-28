#!/usr/bin/env python3
"""Generate SunMRRC SDD V3.4 full architecture PPT — 16 slides, follows SDD chapters exactly."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
import os

# ═══ Theme ═══════════════════════════════════════════════════════
BG      = RGBColor(0x07, 0x0A, 0x12)
BG2     = RGBColor(0x0C, 0x12, 0x20)   # gradient end / panel
CARD    = RGBColor(0x10, 0x16, 0x26)
CARD2   = RGBColor(0x16, 0x1F, 0x34)
INK     = RGBColor(0x0B, 0x10, 0x1C)   # deep ink for shadows
ACCENT  = RGBColor(0x5B, 0x9B, 0xD5)
ACCENT2 = RGBColor(0x3D, 0xD6, 0xC4)   # teal companion for gradients
GREEN   = RGBColor(0x6B, 0xCF, 0x7F)
ORANGE  = RGBColor(0xE8, 0xA8, 0x40)
RED     = RGBColor(0xE0, 0x6A, 0x7A)
PURPLE  = RGBColor(0x9A, 0x7A, 0xE8)
WHITE   = RGBColor(0xEC, 0xF0, 0xF6)
GRAY    = RGBColor(0x80, 0x90, 0xA4)
DIM     = RGBColor(0x44, 0x54, 0x68)
W = Inches(13.333); H = Inches(7.5)
TOTAL = 16

prs = Presentation()
prs.slide_width  = W; prs.slide_height = H

def _shadow(sh, blur=0.10, dist=0.055, alpha=58, dir=5400000):
    """Attach a soft outer drop shadow to a shape via raw OOXML."""
    spPr = sh._element.spPr
    ef = spPr.makeelement(qn('a:effectLst'), {})
    sdw = ef.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(blur*914400)),
        'dist': str(int(dist*914400)),
        'dir': str(dir), 'rotWithShape': '0'})
    clr = sdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    a = clr.makeelement(qn('a:alpha'), {'val': str(alpha*1000)})
    clr.append(a); sdw.append(clr); ef.append(sdw); spPr.append(ef)

def _grad(sh, c1, c2, angle=90):
    """Linear gradient fill on a shape (c1 → c2)."""
    spPr = sh._element.spPr
    for tagn in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
        e = spPr.find(qn(tagn))
        if e is not None: spPr.remove(e)
    g = spPr.makeelement(qn('a:gradFill'), {'flip': 'none', 'rotWithShape': '1'})
    lst = g.makeelement(qn('a:gsLst'), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = lst.makeelement(qn('a:gs'), {'pos': str(pos)})
        sc = gs.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (col[0], col[1], col[2])})
        gs.append(sc); lst.append(gs)
    g.append(lst)
    lin = g.makeelement(qn('a:lin'), {'ang': str(int(angle*60000)), 'scaled': '1'})
    g.append(lin)
    ln = spPr.find(qn('a:ln'))
    if ln is not None: spPr.insert(list(spPr).index(ln), g)
    else: spPr.append(g)

def slide(grad=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    if grad:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.line.fill.background(); _grad(bg, BG, BG2, angle=125)
        bg.shadow.inherit = False
    return s

def txt(s, l, t, w, h, text, sz=12, c=WHITE, b=False, al=PP_ALIGN.LEFT, spacing=None):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.font.name = 'Inter'; p.alignment = al
    if spacing is not None:
        p.font._rPr.set('spc', str(int(spacing*100)))
    return tb

def box(s, l, t, w, h, f=CARD, br=DIM, rad=True, shadow=True, grad=None, lw=0.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    try: sh.adjustments[0] = 0.045
    except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb = f
    sh.line.color.rgb = br; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    if grad: _grad(sh, grad[0], grad[1], grad[2] if len(grad) > 2 else 90)
    if shadow: _shadow(sh)
    return sh

def accent_bar(s, l, t, h, clr, w=0.07):
    """A glowing vertical accent bar — card left-edge highlight."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    try: sh.adjustments[0] = 0.5
    except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background()
    sh.shadow.inherit = False; _shadow(sh, blur=0.08, dist=0, alpha=50)
    return sh

def sep(s, y, w=1.1, clr=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(w), Pt(4))
    try: sh.adjustments[0] = 0.5
    except Exception: pass
    sh.line.fill.background(); sh.shadow.inherit = False
    _grad(sh, ACCENT2, ACCENT, 0); _shadow(sh, blur=0.06, dist=0, alpha=45)

def progress(s, idx):
    """Slim chapter progress bar pinned to the top edge."""
    track = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Pt(4))
    track.fill.solid(); track.fill.fore_color.rgb = CARD2; track.line.fill.background()
    track.shadow.inherit = False
    fillw = int(W * idx / TOTAL)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, fillw, Pt(4))
    bar.line.fill.background(); bar.shadow.inherit = False
    _grad(bar, ACCENT2, ACCENT, 0)

def n(s, num):
    txt(s, 11.9, 7.06, 1.2, 0.3, f'{num:02d} / {TOTAL}', sz=9, c=DIM, al=PP_ALIGN.RIGHT)

def kicker(s, text):
    """Small uppercase eyebrow label above the title."""
    txt(s, 0.82, 0.36, 8, 0.25, text.upper(), sz=10, c=ACCENT2, b=True, spacing=2.5)

def tag(s, l, t, label, clr):
    ln = len(label)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(ln*0.098+0.28), Inches(0.27))
    try: sh.adjustments[0] = 0.5
    except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(clr[0]//4, clr[1]//4, clr[2]//4)
    sh.line.color.rgb = clr; sh.line.width = Pt(0.75); sh.shadow.inherit = False
    tf = sh.text_frame; p = tf.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(7.5); p.font.color.rgb = clr; p.font.bold = True; p.font.name = 'Inter'

def rich(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True; return tb.text_frame

def run(tf, text, sz=11, c=WHITE, b=False):
    p = tf.add_paragraph(); p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.font.name = 'Inter'; return p

def hdr(s, title, sub=None, idx=None, eyebrow=None):
    if idx is not None: progress(s, idx)
    if eyebrow: kicker(s, eyebrow)
    txt(s, 0.8, 0.5, 11.5, 0.55, title, sz=25, b=True)
    sep(s, 1.04)
    if sub: txt(s, 0.82, 1.14, 11.8, 0.32, sub, sz=11, c=GRAY)

def mini_card(s, l, t, w, h, title, desc, tc=ACCENT):
    box(s, l, t, w, h, f=CARD2, br=tc)
    txt(s, l+0.1, t+0.06, w-0.2, 0.28, title, sz=9, c=tc, b=True, al=PP_ALIGN.CENTER)
    txt(s, l+0.1, t+0.34, w-0.2, h-0.42, desc, sz=7.5, c=GRAY, al=PP_ALIGN.CENTER)

def icon_badge(s, l, t, glyph, clr, d=0.5):
    """Rounded glyph badge with tinted fill + glow."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(d), Inches(d))
    try: sh.adjustments[0] = 0.3
    except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(clr[0]//3, clr[1]//3, clr[2]//3)
    sh.line.color.rgb = clr; sh.line.width = Pt(1.0); sh.shadow.inherit = False
    _shadow(sh, blur=0.07, dist=0, alpha=45)
    tf = sh.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.text = glyph; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(int(d*36)); p.font.name = 'Inter'

def arrow(s, x, y): txt(s, x, y, 0.2, 0.2, '→', sz=10, c=DIM, al=PP_ALIGN.CENTER)

def vconn(s, x, y, h, clr=ACCENT, label=None):
    """Vertical gradient connector with a down chevron and optional protocol chip."""
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Pt(2.6), Inches(h))
    sh.line.fill.background(); sh.shadow.inherit = False; _grad(sh, clr, clr, 90)
    tr = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x-0.08), Inches(y+h-0.02), Inches(0.18), Inches(0.13))
    tr.rotation = 180; tr.fill.solid(); tr.fill.fore_color.rgb = clr; tr.line.fill.background()
    tr.shadow.inherit = False
    if label:
        lw = len(label)*0.075 + 0.45
        ch = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.2), Inches(y+h/2-0.16), Inches(lw), Inches(0.32))
        try: ch.adjustments[0] = 0.5
        except Exception: pass
        ch.fill.solid(); ch.fill.fore_color.rgb = INK; ch.line.color.rgb = clr; ch.line.width = Pt(1.0)
        ch.shadow.inherit = False; _shadow(ch, blur=0.06, dist=0.02, alpha=40)
        p = ch.text_frame.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(8.5); p.font.color.rgb = clr; p.font.bold = True; p.font.name = 'Inter'

def chip(s, l, t, w, label, clr, fill=None):
    """A compact rounded component chip with subtle depth."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.34))
    try: sh.adjustments[0] = 0.28
    except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill if fill else RGBColor(clr[0]//4, clr[1]//4, clr[2]//4)
    sh.line.color.rgb = clr; sh.line.width = Pt(0.75); sh.shadow.inherit = False
    _shadow(sh, blur=0.05, dist=0.02, alpha=35)
    p = sh.text_frame.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(8.5); p.font.color.rgb = clr; p.font.bold = True; p.font.name = 'Inter'

def tier_band(s, l, t, h, label, sub, clr):
    """Vertical tier label band on the left edge of a grouped row."""
    sh = box(s, l, t, 1.35, h, f=CARD2, br=clr, lw=1.0)
    _grad(sh, RGBColor(clr[0]//4, clr[1]//4, clr[2]//4), INK, 115)
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12.5); p.font.color.rgb = clr; p.font.bold = True; p.font.name = 'Inter'
    p2 = tf.add_paragraph(); p2.text = sub; p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(7.5); p2.font.color.rgb = GRAY; p2.font.name = 'Inter'

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
s = slide()
# glowing accent rail + radio glyph badge
accent_bar(s, 0.82, 1.55, 2.7, ACCENT2, w=0.09)
icon_badge(s, 11.0, 1.5, '📡', ACCENT, d=1.3)
# eyebrow
txt(s, 1.05, 1.5, 9, 0.3, 'SOFTWARE DESIGN DESCRIPTION  ·  V3.4', sz=12, c=ACCENT2, b=True, spacing=3)
# wordmark
txt(s, 1.0, 1.95, 11, 1.2, 'SunMRRC', sz=66, b=True)
txt(s, 1.05, 3.2, 11, 0.5, 'SunSDR2 DX  ·  Mobile Remote Radio Control', sz=21, c=ACCENT)
sep(s, 3.95, w=1.4)
# feature chip row
feat = [('📡', 'HTTPS / WSS', ACCENT), ('🎛️', 'DSP Pipeline', GREEN),
        ('🎤', 'TX Voice', ORANGE), ('🛡️', 'PTT Safety', RED), ('⚡', 'Per-band Power', PURPLE)]
fx = 1.0
for g, lbl, clr in feat:
    chip(s, fx, 4.25, 2.1, f'{g}  {lbl}', clr)
    fx += 2.25
# stat strip
stats1 = [('16', 'SDD Chapters'), ('12', 'Architecture Decisions'), ('5', 'WebSocket Endpoints'),
          ('8', 'PTT Safety Layers'), ('5', 'SVG Diagrams')]
sx = 1.0
for val, lbl in stats1:
    p = box(s, sx, 4.95, 2.1, 1.0, f=CARD, br=DIM, grad=(CARD2, INK, 120))
    txt(s, sx, 5.12, 2.1, 0.45, val, sz=28, c=ACCENT, b=True, al=PP_ALIGN.CENTER)
    txt(s, sx, 5.6, 2.1, 0.3, lbl, sz=8.5, c=GRAY, al=PP_ALIGN.CENTER)
    sx += 2.25
txt(s, 1.0, 6.25, 11.5, 0.3, 'Python 3.12 · FastAPI / Uvicorn · NumPy / SciPy · Web Audio API · vanilla JavaScript · WDSP', sz=10.5, c=GRAY)
txt(s, 1.0, 6.95, 11.5, 0.3, 'Document ID: SDD-SUNMRRC-2026-001   ·   IBM TeamSD v2.3.2   ·   https://radio.vlsc.net:8889', sz=9.5, c=DIM)
n(s, 1)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary (Ch 1)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Executive Summary', 'SunSDR2 DX mobile-first browser remote control — V1.0 production release',
    idx=1, eyebrow='Chapter 1')
items = [
    ('📱', 'Mobile-first PWA', 'iPhone/Safari primary target · HTTPS/WSS secure context · Touch-optimized · Safe-area support · Service worker · Manifest', ACCENT),
    ('🔊', 'Real-time RX Audio', 'IQ demodulation (USB/LSB/AM/FM) · Tagged dual-codec (Opus 28kbps / Int16 PCM) · WDSP NR2/AGC · S-meter · Waterfall', GREEN),
    ('🎤', 'TX Voice Modulation', 'Web Audio EQ chain (4 presets) · AudioWorklet 48k→16k · Opus encode · Hilbert SSB · tanh soft limiter · DRIVE power control · On-air verified', ORANGE),
    ('⚡', 'Per-band TX Power', 'Device DRIVE (0x0017) sqrt-taper · /api/band_power + Band Power UI · band_power.json persistence · Re-sent on QSY/PTT', PURPLE),
    ('🛡️', 'PTT Safety System', '8-layer defense in depth · ACK retry 3× · Backup s: channel · 30s watchdog · Touch-move-out release · Forced-RX handler', RED),
    ('📐', 'Architecture Record', '16 SDD chapters · 12 ADs · 5 SVG diagrams · 16-slide PPTX · TeamSD methodology · Full protocol documentation', ACCENT2),
]
y = 1.72
for icon, title, desc, clr in items:
    box(s, 0.5, y, 12.33, 0.82, f=CARD, br=RGBColor(clr[0]//3, clr[1]//3, clr[2]//3), lw=0.75,
        grad=(CARD, INK, 115))
    accent_bar(s, 0.5, y+0.1, 0.62, clr)
    icon_badge(s, 0.78, y+0.16, icon, clr, d=0.5)
    txt(s, 1.48, y+0.11, 5, 0.28, title, sz=13.5, c=clr, b=True)
    txt(s, 1.48, y+0.42, 11.0, 0.36, desc, sz=9.5, c=GRAY)
    y += 0.88
n(s, 2)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Business Direction + Project Definition (Ch 2+3)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Business Direction · Project Definition', 'Goals · Objectives · Scope · Success Criteria', idx=3, eyebrow='Chapters 2–3')

box(s, 0.4, 1.3, 6.2, 5.7, f=CARD, br=ACCENT)
txt(s, 0.6, 1.35, 5.5, 0.3, '🎯  Business Goals & Objectives', sz=15, c=ACCENT, b=True)
goals = [
    ('G1', 'Mobile RX Confidence', GREEN),
    ('G2', 'Safe Remote Control', GREEN),
    ('G3', 'Minimal Deployment', GREEN),
    ('G4', 'Design Continuity', GREEN),
    ('G5', 'Incremental Extensibility', ORANGE),
]
gy = 1.8
for gid, gdesc, gc in goals:
    tag(s, 0.7, gy, gid, gc)
    txt(s, 1.3, gy, 4.5, 0.25, gdesc, sz=11, c=WHITE)
    gy += 0.35
txt(s, 0.7, gy+0.1, 5.5, 0.22, '8 Objectives — 6 Implemented · 1 Open (ATR)', sz=10, c=DIM)
txt(s, 0.7, gy+0.4, 5.5, 0.22, '5 Strategies — Mobile-first · Direct UDP · Browser audio', sz=10, c=DIM)
txt(s, 0.7, gy+0.7, 5.5, 0.22, '6 Tactics — HTTPS default · WSS select · PTT ACK loop', sz=10, c=DIM)

box(s, 6.9, 1.3, 6.0, 5.7, f=CARD, br=GREEN)
txt(s, 7.1, 1.35, 5.5, 0.3, '📋  Project Scope', sz=15, c=GREEN, b=True)

def place_text(s, l, t, w, items, color):
    for i, item in enumerate(items):
        txt(s, l, t+i*0.19, w, 0.18, item, sz=9, c=color)

txt(s, 7.1, 1.75, 5.5, 0.2, '▶  In Scope', sz=12, c=GREEN, b=True)
place_text(s, 7.3, 2.05, 5.3, [
    '  • Mobile web UI', '  • 5 WebSocket endpoints',
    '  • RX audio + waterfall', '  • TX mic → Hilbert SSB → IQ',
    '  • Per-band DRIVE power', '  • /api/band_power + mem_channels',
    '  • HTTPS/WSS auto-start', '  • Direct SunSDR2 DX UDP',
    '  • DSP mode + WDSP toggles', '  • SDD documentation',
], GRAY)
txt(s, 7.1, 4.05, 5.5, 0.2, '⊘  Out of Scope', sz=12, c=RED, b=True)
place_text(s, 7.3, 4.35, 5.3, ['  • Native iOS/Android app',
    '  • Cloud multi-tenant', '  • ATR-1000 backend (stub only)',
    '  • CW/FT8/Logbook pages', '  • Hamlib rig abstraction',
], DIM)
txt(s, 7.1, 5.4, 5.5, 0.2, '✅  9 Success Criteria — All Met', sz=11, c=GREEN, b=True)
n(s, 3)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — System Context (Ch 4)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'System Context', 'Actors · External Interfaces · Data Flows · System Boundaries', idx=4, eyebrow='Chapter 4 · APP 011')

layers = [
    (1.4, '📱  Mobile Browser', 'iOS Safari / Chrome · Web Audio API · WASM Opus\nTouch UX · PWA Manifest · Service Worker · Canvas', ACCENT),
    (2.6, 'HTTPS/WSS  —  5 WebSocket Endpoints', '/WSCTRX (control) · /WSaudioRX (RX audio) · /WSaudioTX (TX mic)\n/WSspectrum (waterfall) · /WSATR1000 (tuner placeholder)', GREEN),
    (3.8, '🖥️  FastAPI SunMRRC Server', 'server.py · Static files · Auth middleware · TLS auto-detect\nIQ processing loop · TX pacer thread · /api/band_power + mem_channels', ORANGE),
    (5.0, 'imports  →  web_control/', 'sunsdr_direct.py (UDP protocol) · dsp.py (IQ DSP + TX modulator)\nopus_rx.py (server-side Opus) · wdsp_wrapper.py (NR2/AGC)', PURPLE),
    (6.2, '📡  SunSDR2 DX Hardware', 'UDP :50001 Control (DRIVE 0x0017 · PTT · Frequency)\nUDP :50002 IQ Stream (0xFFFE RX · 0xFFFD TX · 0x1F00 Telemetry)', RED),
]
for y, title, desc, color in layers:
    box(s, 0.5, y, 12.3, 1.0, f=CARD, br=color)
    txt(s, 0.7, y+0.06, 11.5, 0.28, title, sz=13, c=color, b=True)
    txt(s, 0.7, y+0.35, 11.5, 0.6, desc, sz=10, c=GRAY)
    if 'HTTPS' not in title and 'imports' not in title and 'Hardware' not in title:
        pass

txt(s, 0.8, 7.15, 11, 0.25, '4 Actors: HAM Operator · System Maintainer · Browser Runtime · SunSDR2 DX   |   11 External Interfaces   |   9 Data Flows', sz=9, c=DIM)
n(s, 4)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — Non-Functional Requirements (Ch 5)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Non-Functional Requirements', '7 categories · 25 NFRs — Performance · Security · Compatibility · Operability · Audio Quality', idx=5, eyebrow='Chapter 5 · ART 0507')

cats = [
    ('⚡  Performance', '7 NFRs', [
        'NFR-001  RX audio latency: low enough for live monitoring',
        'NFR-002  Control response <200ms · PING→PONG ms display',
        'NFR-003  Spectrum ~512B/frame @ 38Hz (~19 KB/s)',
        'NFR-004  Audio Opus ~18-24kbps (PCM ~256kbps)',
        'NFR-005  CPU stability: no sustained DSP overload',
        'NFR-006  Client fan-out: multi-browser without crash',
        'NFR-007  Waterfall: adaptive noise floor, smooth render',
    ], ACCENT),
    ('🔒  Security', '4 NFRs', [
        'NFR-020  HTTPS/WSS secure origin for mobile production',
        'NFR-021  TLS key files never embedded in docs/code',
        'NFR-022  HTTP fallback only when certs missing or DISABLE_SSL=1',
        'NFR-023  Auth boundary: no false claim of server-side auth',
    ], RED),
    ('🎵  Audio Quality', '4 NFRs', [
        'NFR-060  RX Int16 PCM → Float32 decode path',
        'NFR-061  USB/LSB/AM/FM/CW demodulation quality',
        'NFR-062  WDSP NR2/NB/ANF/NF/AGC when libwdsp loads',
        'NFR-063  TX: Hilbert SSB + tanh soft limiter, clean envelope',
        'NFR-065  Per-band DRIVE (0x0017) sets TX power',
    ], GREEN),
    ('🔧  Operability + Maintainability', '6 NFRs', [
        'NFR-040  Logging: startup, TLS, radio connect, IQ bind',
        'NFR-041  Env config: DEVICE_HOST, WEB_PORT, DISABLE_SSL',
        'NFR-042  Certificate expiry monitoring scripts',
        'NFR-043  Service worker bypasses JS/HTML cache',
        'NFR-050  Small backend: handlers understandable in server.py',
        'NFR-051  Explicit gaps: planned hooks documented until removed',
    ], ORANGE),
]
# 2×2 grid layout
positions = [(0.5, 1.4), (6.8, 1.4), (0.5, 4.35), (6.8, 4.35)]
for (cat_title, count, items, clr), (px, py) in zip(cats, positions):
    box(s, px, py, 6.05, 2.8, f=CARD, br=clr)
    txt(s, px+0.2, py+0.08, 5.5, 0.22, f'{cat_title}  ({count})', sz=12, c=clr, b=True)
    iy = py+0.45
    for item in items:
        txt(s, px+0.3, iy, 5.6, 0.2, item, sz=8.5, c=GRAY)
        iy += 0.235
n(s, 5)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — Use Case Model (Ch 6)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Use Case Model', '4 Actors · 6 Core Use Cases · 6 Extended / Planned Use Cases', idx=6, eyebrow='Chapter 6 · ART 0508')

ucs = [
    ('UC-001', 'Start Mobile Session', 'Open https→login→power on→WS connections→UI live', '✅', GREEN),
    ('UC-002', 'Remote Receive Audio', 'IQ→decode→demod→Opus/PCM→/WSaudioRX→play', '✅', GREEN),
    ('UC-003', 'Tune Frequency & Mode', 'setFreq/setMode→radio/DSP→getFreq/getMode ack', '✅', GREEN),
    ('UC-004', 'Monitor Spectrum & S-Meter', 'FFT→quantize→/WSspectrum→adaptive floor+colour ramp', '✅', GREEN),
    ('UC-005', 'PTT & Tune Control', 'PTT press→TX→release→ACK retry→backup s:→watchdog', '✅', GREEN),
    ('UC-006', 'Adjust DSP & Audio', 'WDSP/NR2/NB/AGC/gain/filter→DSP setters→broadcast', '✅', GREEN),
    ('UC-007', 'Browser Mic TX Audio', 'Mic→EQ→AudioWorklet→Opus→WS→Hilbert SSB→0xFFFD IQ', '✅', GREEN),
    ('UC-008', 'Memory Channel Save/Recall', 'Menu→edit→POST /api/mem_channels→mem_channels.json', '✅', GREEN),
    ('UC-013', 'Per-band TX Power', 'Band Power menu→edit drive %→POST /api/band_power→apply', '✅', GREEN),
    ('UC-009', 'ATR-1000 Power/SWR', 'Website placeholders; /WSATR1000 stub only', '⚠', ORANGE),
    ('UC-010/11', 'CW Decoder / FT8 Automation', 'Menu links removed (pages absent)', '✕', RED),
    ('UC-012', 'Recordings', 'recordings.html + server-side RX MP3 via ffmpeg', '✅', GREEN),
]
y = 1.5
for uid, name, desc, status, sc in ucs:
    box(s, 0.5, y, 12.3, 0.43, f=CARD, br=sc)
    tag(s, 0.65, y+0.07, uid, sc)
    txt(s, 1.55, y+0.06, 3.5, 0.22, name, sz=11, c=WHITE, b=True)
    txt(s, 5.1, y+0.06, 6.5, 0.22, desc, sz=9, c=GRAY)
    txt(s, 11.5, y+0.1, 1, 0.2, status, sz=16, c=sc, b=True, al=PP_ALIGN.CENTER)
    y += 0.48

n(s, 6)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Subject Area Model (Ch 7)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Subject Area Model', '7 Subject Areas · 18 Entity Definitions · State Ownership', idx=7, eyebrow='Chapter 7 · APP 408')

areas = [
    ('ClientSession', 'WS memberships, UI state', 'Browser cookies/JS'),
    ('RadioControl', 'Freq · PTT · Gain · Filter · Tune', 'Runtime (SunSDR2DXClient)'),
    ('IQStream', 'UDP intake → complex samples → DSP feed', 'In-process (UDP socket)'),
    ('RXAudioFlow', 'Demodulation → Opus/PCM → playback', 'In-process + WS broadcast'),
    ('SpectrumFlow', 'FFT · S-meter · Waterfall render state', 'In-process + WS broadcast'),
    ('DSPConfig', 'Mode · WDSP · AGC · Notches · Filters', 'Runtime + cookie replay'),
    ('OperationalConfig', 'Host · Port · TLS · Env vars · Scripts', 'Filesystem (env, certs, logs)'),
]
y = 1.5
for area, desc, persist in areas:
    box(s, 0.5, y, 12.3, 0.71, f=CARD, br=ACCENT)
    txt(s, 0.7, y+0.06, 3, 0.22, area, sz=12, c=ACCENT, b=True)
    txt(s, 0.7, y+0.32, 5.5, 0.28, desc, sz=10, c=GRAY)
    txt(s, 6.5, y+0.06, 6, 0.22, persist, sz=9, c=DIM)
    y += 0.78

# Key entities
txt(s, 0.8, 7.0, 11, 0.22, 'Key entities: TXModulationFrame · TXDriveConfig · BandPowerConfig · MemoryChannel · WaterfallRenderState · SMeterState · WDSPConfig · TLSConfig', sz=9, c=DIM)
n(s, 7)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Architecture Decisions (Ch 8)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Architecture Decisions', '12 ADs — Architectural · Design · Performance · Safety · Extensibility', idx=8, eyebrow='Chapter 8 · ART 0513')

ads = [
    ('AD-001', 'FastAPI/Uvicorn', 'Async server, native WS, small surface', GREEN),
    ('AD-002', 'Direct UDP Control', 'Verified SunSDR protocol, no Hamlib', GREEN),
    ('AD-003', 'DSP-Owned Mode', 'Hardware mode-agnostic, SW demod selects', GREEN),
    ('AD-004', 'Tagged Dual-Codec Audio', 'Opus+PCM RX and TX, 1-byte tag, no race', GREEN),
    ('AD-005', 'Quantized Spectrum', '512 uint8 bytes/frame, compact waterfall', GREEN),
    ('AD-006', 'HTTPS/WSS Default', 'iOS secure context, mic+AudioContext require', GREEN),
    ('AD-007', 'PTT Release Safety', '8-layer defense in depth, ACK retry, backup s:', GREEN),
    ('AD-008', 'Optional WDSP', 'Falls back cleanly, UI queries availability', GREEN),
    ('AD-009', 'Gap Tracking', 'Document hooks until backend exists or removed', ORANGE),
    ('AD-010', 'Device DRIVE Power', '0x0017 per-band, sqrt taper byte, re-send QSY', GREEN),
    ('AD-011', '0x1F00 Telemetry', 'off30 f32=W · off16 u16/10=V · off18=°C · NO SWR from device', GREEN),
    ('AD-012', 'TX Gain Staging', 'Preamp ×1.5 + drive gain ×3.0 + tanh@1.0 ≈ 4% engagement', GREEN),
]
y = 1.35
for adid, topic, desc, clr in ads:
    box(s, 0.5, y, 12.3, 0.44, f=CARD, br=DIM)
    tag(s, 0.65, y+0.08, adid, clr)
    txt(s, 1.55, y+0.07, 4, 0.3, topic, sz=11, c=WHITE, b=True)
    txt(s, 5.8, y+0.07, 6.5, 0.3, desc, sz=10, c=GRAY)
    y += 0.49

n(s, 8)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Architecture Overview: visual block diagram (Ch 9)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Architecture Overview', 'Four-layer stack — Browser → FastAPI server → shared DSP → SunSDR2 DX hardware', idx=9, eyebrow='Chapter 9 · ART 0512')

# ── Layer 1: Browser ──
box(s, 1.2, 1.35, 10.9, 1.05, f=CARD, br=ACCENT)
txt(s, 1.4, 1.4, 6, 0.26, '📱  Mobile Browser  (iOS Safari / Chrome)', sz=13, c=ACCENT, b=True)
cw = 2.05
for i, (lbl) in enumerate(['controls.js', 'mobile.js', 'tx_audio_eq.js', 'ptt_manager.js', 'WASM Opus']):
    chip(s, 1.4 + i*2.14, 1.78, cw, lbl, ACCENT)
txt(s, 1.4, 2.13, 10.5, 0.24, 'Web Audio API · AudioWorklet 48k→16k · Canvas waterfall · Touch PTT · PWA / Service Worker', sz=8.5, c=GRAY)

# ── connector with WS endpoints ──
vconn(s, 6.6, 2.42, 0.55, ACCENT, 'HTTPS / WSS')
txt(s, 7.4, 2.46, 5.6, 0.5, '/WSCTRX · /WSaudioRX · /WSaudioTX\n/WSspectrum · /WSATR1000   (?token= auth)', sz=8.5, c=DIM)

# ── Layer 2: FastAPI server ──
box(s, 1.2, 3.05, 10.9, 1.15, f=CARD, br=GREEN)
txt(s, 1.4, 3.1, 6, 0.26, '🖥️  FastAPI / Uvicorn Server  (server.py)', sz=13, c=GREEN, b=True)
for i, lbl in enumerate(['Auth + TLS', 'IQ Loop', 'TX Pacer', '/api/band_power', '/api/mem_channels']):
    chip(s, 1.4 + i*2.14, 3.48, cw, lbl, GREEN)
txt(s, 1.4, 3.84, 10.5, 0.3, 'Static serving · 5 WS endpoints · client fan-out · IQ decode → DSP feed · adaptive TX pacing (5.12ms/pkt) · 0x1F00 telemetry', sz=8.5, c=GRAY)

# ── connector ──
vconn(s, 6.6, 4.22, 0.5, PURPLE, 'in-process import')

# ── Layer 3: shared DSP ──
box(s, 1.2, 4.8, 10.9, 1.05, f=CARD, br=PURPLE)
txt(s, 1.4, 4.85, 7, 0.26, '⚙️  Shared DSP / Protocol  (web_control/)', sz=13, c=PURPLE, b=True)
for i, lbl in enumerate(['sunsdr_direct', 'dsp.py', 'TXModulator', 'opus_rx', 'wdsp_wrapper']):
    chip(s, 1.4 + i*2.14, 5.23, cw, lbl, PURPLE)
txt(s, 1.4, 5.58, 10.5, 0.24, 'UDP protocol · SpectrumProcessor (FFT) · AudioDemodulator (SSB/AM/FM) · Hilbert SSB + tanh · server Opus · NR2/AGC', sz=8.5, c=GRAY)

# ── connector ──
vconn(s, 6.6, 5.87, 0.5, RED, 'UDP :50001 / :50002')

# ── Layer 4: hardware ──
box(s, 1.2, 6.45, 10.9, 0.7, f=CARD, br=RED)
txt(s, 1.4, 6.5, 7, 0.26, '📡  SunSDR2 DX Hardware  (192.168.16.200)', sz=13, c=RED, b=True)
txt(s, 1.4, 6.82, 10.5, 0.24, 'Ctrl :50001 (DRIVE 0x0017 · PTT · Freq)   ·   IQ :50002 (0xFFFE RX · 0xFFFD TX · 0x1F00 telem)   ·   PA — no ALC', sz=8.5, c=GRAY)
n(s, 9)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Architecture Overview: RX + TX Signal Chains (Ch 9)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Signal Processing — RX + TX Chains', 'IQ demodulation downlink · Hilbert SSB modulation uplink · gain staging', idx=10, eyebrow='Chapter 9 · Signal Flow')

# RX chain
box(s, 0.4, 1.3, 12.5, 2.7, f=CARD, br=GREEN)
txt(s, 0.6, 1.35, 6, 0.28, '🔊  RX: IQ → Audio', sz=14, c=GREEN, b=True)
rx_stages = ['UDP\nIntake', '24-bit\nDecode', 'IF Shift\n+30.5kHz', 'SSB\nBandpass', '5:1\nDecimate', 'AGC\n(WDSP opt)', 'Opus\nEncode', '/WSaudioRX\nBroadcast', 'Browser\nPlayback']
rxx = 0.6
for stage in rx_stages:
    box(s, rxx, 1.75, 1.25, 0.8, f=CARD2, br=GREEN)
    txt(s, rxx+0.08, 1.82, 1.1, 0.65, stage, sz=8, c=GREEN, b=True, al=PP_ALIGN.CENTER)
    if rxx > 0.6: arrow(s, rxx-0.22, 2.0)
    rxx += 1.35
txt(s, 0.6, 2.7, 12, 0.5, 'Waterfall: IQ→FFT(2048,hanning)→dB clip[-120,0]→uint8 quantize 512B→/WSspectrum @~38Hz→⏫ accumulate 10→30th %ile noise floor→blue bias+contrast gain→colour ramp', sz=9, c=DIM)
txt(s, 0.6, 3.0, 12, 0.5, 'S-Meter: FFT→90th percentile→S9=9+(p90+73)/6→asymmetric EMA(attack α=0.5,release α=0.15)→stable needle,fast-rise slow-decay', sz=9, c=DIM)
txt(s, 0.6, 3.3, 12, 0.4, 'Codec: Server ctypes libopus · bitrate via opus_encode() max_data_bytes cap (avoids arm64 variadic ctl) · Client WASM OpusDecoder · 1-byte tag per frame · Default Opus ~18-24kbps', sz=9, c=DIM)

# TX chain
box(s, 0.4, 4.25, 12.5, 2.85, f=CARD, br=ORANGE)
txt(s, 0.6, 4.3, 6, 0.28, '🎤  TX: Mic → RF (Client + Server)', sz=14, c=ORANGE, b=True)
txt(s, 0.6, 4.65, 12, 0.2, '📱  Client: micSource → preamp(×1.5) → antiAlias(LP 4.5kHz ×2) → eqLow(350Hz pk) → eqMid(1500Hz pk) → eqHigh(2700Hz sh) → midCut → presence → compressor(3:1,thr=-24) → makeup(×1.6) → noiseGate → AudioWorklet(48k→16k Int16) → Opus encode → /WSaudioTX', sz=9, c=GRAY)
txt(s, 0.6, 5.1, 12, 0.2, '🖥️  Server: TxOpusDecoder → TXModulator.feed_audio() → fractional resampler(16k→15625Hz) → overlap-save Hilbert SSB(MARGIN=256,80-sample hops) → upsample ×2.5(→39063Hz) → drive gain(×3.0×drive%) → tanh(1.0) soft limiter → 24-bit IQ packing(vectorized numpy) → jitter buffer(prime 24/reprime 12) → TX pacer(5.12ms/pkt,adaptive ±15%) → 0xFFFD UDP :50002', sz=9, c=GRAY)

txt(s, 0.8, 5.55, 12, 0.2, '📊  Gain Staging (AD-012):  in peak ~0.50 → Hilbert peak ~0.65 → drive gain peak ~1.95 → tanh(1.0) peak ~0.96 (4% reduction)', sz=9, c=GREEN, b=True)
txt(s, 0.8, 5.8, 12, 0.2, '❌  OLD (preamp ×3.0): in peak=1.00 → Hilbert peak=1.32 → drive gain peak=3.95 → tanh peak=0.999 (75% reduction → HEAVY DISTORTION on every loud syllable)', sz=9, c=RED)

txt(s, 0.6, 6.15, 12, 0.2, '📡  Device: 0xFFFD @195 pkt/s (39063 Hz IQ) → DRIVE command (0x0017) per-band power → PA (no ALC) → Antenna  |  ⚡ Tune bypasses entire chain: pre-computed IQ, TX_TUNE_SCALE=0.35 → ~10W safe continuous', sz=9, c=DIM)
txt(s, 0.6, 6.45, 12, 0.2, '🔑  RF power = DEVICE DRIVE only.  Client gain = healthy digital level.  Server gain = lifts into tanh knee.  Turn drive ↑ for more power, not client gain.', sz=10, c=WHITE, b=True)
n(s, 10)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Service + Component Model (Ch 10+11)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Service & Component Model', '12 services · backend / frontend / shared DSP components', idx=11, eyebrow='Chapters 10–11')

# Left: Services
box(s, 0.4, 1.3, 6.2, 5.8, f=CARD, br=GREEN)
txt(s, 0.6, 1.35, 5.5, 0.28, '🔌  Service Portfolio', sz=15, c=GREEN, b=True)
svcs = [
    ('StaticUIService', 'mobile PWA assets from static/'),
    ('ControlService', '/WSCTRX command dispatch + state sync'),
    ('RXAudioService', 'IQ→PCM→Opus /WSaudioRX fan-out'),
    ('TXAudioIngressService', '/WSaudioTX mic→TXModulator'),
    ('SpectrumService', 'FFT quantize /WSspectrum fan-out'),
    ('SunSDRDeviceService', 'UDP protocol, boot, heartbeat'),
    ('IQProcessingService', 'receive, decode, feed DSP'),
    ('DSPService', 'demodulation + optional WDSP chain'),
    ('TXPowerService', '0x0017 DRIVE per-band + /api/band_power'),
    ('BandPowerService', 'per-band drive % JSON persistence'),
    ('MemoryChannelService', '/api/mem_channels JSON persistence'),
    ('TLSStartupService', 'cert/key auto-detect, HTTP fallback'),
]
sy = 1.75
for name, desc in svcs:
    txt(s, 0.7, sy, 2.8, 0.22, name, sz=9, c=GREEN, b=True)
    txt(s, 3.6, sy, 2.8, 0.22, desc, sz=8, c=GRAY)
    sy += 0.3

# Right: Components
box(s, 6.9, 1.3, 6.0, 5.8, f=CARD, br=ACCENT)
txt(s, 7.1, 1.35, 5.5, 0.28, '🧩  Key Components', sz=15, c=ACCENT, b=True)
txt(s, 7.1, 1.75, 5.5, 0.2, 'Backend (Python)', sz=11, c=ACCENT, b=True)
comps_be = ['FastAPIApp — routes, startup, static serving',
    'RadioClient — UDP protocol (sunsdr_direct.py)',
    'IQLoop — UDP socket, packet validation, DSP feed',
    'TXModulator — Hilbert SSB, drive gain, tanh, IQ pack',
    'TXPacer — adaptive pacing, jitter buffer',
    'RxOpusEncoder / TxOpusDecoder — ctypes libopus',
    'BandPowerAPI / MemChannelAPI — JSON persistence']
cy = 2.05
for c in comps_be:
    txt(s, 7.3, cy, 5.3, 0.2, f'• {c}', sz=8, c=GRAY)
    cy += 0.22

txt(s, 7.1, 3.8, 5.5, 0.2, 'Frontend (JavaScript)', sz=11, c=ACCENT, b=True)
comps_fe = ['controls.js — WS, codec, waterfall, S-meter',
    'mobile.js — UX, menus, DSP panel, settings',
    'tx_audio_eq.js — 4 EQ presets, compressor, preamp',
    'tx_capture_worklet.js — AudioWorklet 48k→16k',
    'ptt_manager.js — ACK retry, watchdog, state sync',
    'tx_button.js — touch PTT, lock, warm-up frames',
    'opus_codec.js — WASM Opus encode/decode (28kbps)']
cy = 4.1
for c in comps_fe:
    txt(s, 7.3, cy, 5.3, 0.2, f'• {c}', sz=8, c=GRAY)
    cy += 0.22

txt(s, 7.1, 5.9, 5.5, 0.2, 'Shared DSP (Python)', sz=11, c=ACCENT, b=True)
txt(s, 7.3, 6.2, 5.3, 0.35, '• StreamProcessor · SpectrumProcessor · AudioDemodulator · TXModulator · Opus Codec · WDSPWrapper', sz=8, c=GRAY)
n(s, 11)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Operational Model (Ch 12)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Operational Model', 'Startup · Configuration · TLS · Connection Matrix · Verification', idx=12, eyebrow='Chapter 12 · ART 0522')

box(s, 0.4, 1.3, 6.2, 3.0, f=CARD, br=ACCENT)
txt(s, 0.6, 1.35, 5.5, 0.25, '🚀  Startup & Config', sz=14, c=ACCENT, b=True)
txt(s, 0.6, 1.7, 5.8, 2.3, '''restart.sh → kill old by cwd → clear port → background
start.sh → simple foreground launcher
restart.sh -f → foreground with live logs
DISABLE_SSL=1 → HTTP debug (no iOS mic support)

Environment:
  DEVICE_HOST=192.168.16.200 (default)
  WEB_PORT=8889 (default), restart.sh uses 8080
  WEB_PASSWORD=sunmrrc (override with env)

Runtime persistence:
  band_power.json → per-band drive % (Band Power UI)
  mem_channels.json → memory channels (/api/mem_channels)
  server.log → runtime logs (when background-started)''', sz=9, c=GRAY)

box(s, 6.9, 1.3, 6.0, 3.0, f=CARD, br=GREEN)
txt(s, 7.1, 1.35, 5.5, 0.25, '🔐  TLS + Auth', sz=14, c=GREEN, b=True)
txt(s, 7.1, 1.7, 5.8, 2.3, '''TLS Auto-detection:
  certs/fullchain.pem  +  certs/radio.vlsc.net.key
  Both present → HTTPS/WSS (iOS requires)
  Missing or DISABLE_SSL=1 → HTTP fallback
  check_ssl_expiry.sh for cert monitoring

Auth (password-based session):
  /login → POST /api/auth/login → sunmrrc_auth cookie
  Token lifetime: 30 days (cookie max-age)
  All routes + WS require ?token= query param
  _auth_tokens server-side set, invalidated on restart

Connection Matrix:
  Browser ↔ SunMRRC : HTTPS/WSS
  SunMRRC ↔ SunSDR2 DX : UDP :50001 / :50002
  IPs: 192.168.16.100 ← → 192.168.16.200''', sz=9, c=GRAY)

box(s, 0.4, 4.55, 12.5, 2.6, f=CARD, br=ORANGE)
txt(s, 0.6, 4.6, 6, 0.25, '🔍  Verification Checklist', sz=14, c=ORANGE, b=True)
checks = [
    ('Verify HTTPS', 'server.log → "sunmrrc https://..."'),
    ('Verify Radio', 'server.log → "SunSDR2DX: True"'),
    ('Verify RX', 'UI /WSaudioRX connected, waterfall flowing'),
    ('Verify Codec', 'Opus active → bitrate ~18-24kbps'),
    ('Verify PTT Safety', 'getPTT:false after release, release ACK confirmed'),
    ('Verify TX Power', 'W= in server.log from 0x1F00 off30 f32 forward watts'),
    ('Verify TX Quality', 'TX chain lim: peak ~0.96 (not 0.999 → distorted!)'),
    ('Verify Gain Staging', 'in peak ~0.5, drv peak ~2.0, lim peak ~0.96 @ 1 Hz'),
    ('Diagnose Jitter', '/tmp/tx_probe.csv + /tmp/tx_rx_probe.csv'),
    ('Level Probes', 'TX chain in:/an:/drv:/lim: lines at 1 Hz in server.log'),
]
cx, cy_val = [0.6, 6.8], [5.0, 5.0]
ci = 0
for check_name, check_desc in checks:
    x = cx[ci % 2]
    txt(s, x, cy_val[ci % 2], 2.5, 0.2, check_name, sz=9.5, c=ORANGE, b=True)
    txt(s, x+2.6, cy_val[ci % 2], 3.8, 0.2, check_desc, sz=9, c=GRAY)
    cy_val[ci % 2] += 0.22
    ci += 1
n(s, 12)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — Feasibility + Known Gaps (Ch 13)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Feasibility Assessment & Known Gaps', 'Production-ready feasibility · open issues · risk matrix', idx=13, eyebrow='Chapter 13 · ART 0530')

box(s, 0.4, 1.3, 6.2, 3.0, f=CARD, br=GREEN)
txt(s, 0.6, 1.35, 5.5, 0.25, '✅  Feasibility: Production Ready', sz=14, c=GREEN, b=True)
feas = [
    ('RX Technical', 'High', 'IQ→demod→broadcast→playback — implemented'),
    ('Control', 'High', 'All radio + DSP commands implemented'),
    ('TX Voice', 'High', 'Full pipeline: mic→Hilbert SSB→IQ, on-air verified'),
    ('Mobile', 'High', 'HTTPS/WSS solves iOS blocker, all features on phone'),
    ('Operational', 'High', 'Single process, restart script, simple ops'),
    ('Extensibility', 'Medium-High', 'WDSP optional, ATR hook preserved, clean module boundaries'),
]
fy = 1.8
for dim, level, expl in feas:
    lc = GREEN if level == 'High' else ORANGE
    txt(s, 0.7, fy, 2, 0.2, dim, sz=10, c=WHITE, b=True)
    tag(s, 2.8, fy+0.02, level, lc)
    txt(s, 4.2, fy, 2.5, 0.2, expl, sz=9, c=GRAY)
    fy += 0.33

box(s, 0.4, 4.5, 6.2, 2.7, f=CARD, br=RED)
txt(s, 0.6, 4.55, 5.5, 0.25, '⚠  Open Issues', sz=14, c=RED, b=True)
txt(s, 0.6, 4.9, 5.8, 2.0, '''I3: /WSATR1000 missing — stub only, no real HW interface
     ATR-1000 is the only SWR source (device has no reverse-power)

I5: Fixed LAN IPs (192.168.16.100 / 192.168.16.200)
     Should move to env/config for deployment flexibility

I7: start.sh prints HTTP URL even with TLS present
     Align script message with server TLS behavior

Resolved (2026-06-25/26):
  I10: TX power formula — now off30 f32 direct float watts
  I11: TX voice distortion — client preamp 3.0→1.5
  I1-I2, I4, I6, I8, I9: all resolved in prior releases''', sz=9, c=GRAY)

box(s, 6.9, 1.3, 6.0, 5.9, f=CARD, br=RED)
txt(s, 7.1, 1.35, 5.5, 0.25, '📊  Risk Matrix', sz=14, c=RED, b=True)
risks = [
    ('R1', 'iOS HTTP instead of HTTPS', 'Medium', RED),
    ('R2', 'Fixed LAN IP mismatch', 'Medium', ORANGE),
    ('R3', 'TX release lost (half-open socket)', 'Low', GREEN),
    ('R4', 'WDSP library unavailable', 'Low', GREEN),
    ('R8', 'Stale frontend assets', 'Low', GREEN),
    ('R9', 'Certificate expiry', 'Low', GREEN),
    ('R10', 'PA over-drive (no ALC)', 'Low', GREEN),
]
ry = 1.75
for rid, rdesc, rprob, rc in risks:
    box(s, 7.1, ry, 5.6, 0.42, f=CARD2, br=rc)
    tag(s, 7.2, ry+0.08, rid, rc)
    txt(s, 7.8, ry+0.06, 3, 0.22, rdesc, sz=10, c=WHITE)
    tag(s, 11.2, ry+0.08, rprob, rc)
    ry += 0.5

txt(s, 7.1, 5.2, 5.5, 0.2, 'Resolved Risks', sz=13, c=GREEN, b=True)
txt(s, 7.1, 5.5, 5.5, 1.5, 'R5: TX audio not transmitted ✓\nR6: ATR frontend connection errors (mitigated) ✓\nR7: Memory channel API missing ✓\nI11: TX voice distortion (gain staging) ✓\nI10: TX power cubic fit wrong ✓', sz=9, c=GRAY)
n(s, 13)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — PTT Safety Architecture (Ch 15)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'PTT Safety — Defense in Depth', 'Release is more safety-critical than keying. No single point of failure can trap the radio in TX.', idx=14, eyebrow='Chapter 15 · Safety')

# Three tiers, each a horizontal band: tier label + its layer cards side by side.
def ptt_layer(s, l, t, w, h, num, title, clr, lines):
    box(s, l, t, w, h, f=CARD, br=clr)
    # number badge
    bd = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l+0.13), Inches(t+0.13), Inches(0.34), Inches(0.34))
    bd.fill.solid(); bd.fill.fore_color.rgb = clr; bd.line.fill.background()
    p = bd.text_frame.paragraphs[0]; p.text = str(num); p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13); p.font.color.rgb = BG; p.font.bold = True; p.font.name = 'Inter'
    txt(s, l+0.56, t+0.11, w-0.7, 0.26, title, sz=11.5, c=clr, b=True)
    tf = rich(s, l+0.16, t+0.46, w-0.3, h-0.5)
    for ln in lines:
        run(tf, ln, sz=8, c=GRAY)

TIER_X = 1.75
TIER_W = 11.1
CARD_GAP = 0.18

# ── Tier A: LOCAL (layers 1-3) ──
tier_band(s, 0.35, 1.35, 1.65, 'LOCAL', 'browser-side', GREEN)
cw3 = (TIER_W - 2*CARD_GAP) / 3
ptt_layer(s, TIER_X, 1.35, cw3, 1.65, 1, 'Touch UX', GREEN,
          ['touchend · move-out', 'cancel · mouseleave', 'touchId tracking', 'tx_button.js'])
ptt_layer(s, TIER_X+cw3+CARD_GAP, 1.35, cw3, 1.65, 2, 'State Machine', GREEN,
          ['isProcessing lock', 'pendingStop queue', '3s leak → unlock', 'tx_button.js'])
ptt_layer(s, TIER_X+2*(cw3+CARD_GAP), 1.35, cw3, 1.65, 3, 'Watchdog', ORANGE,
          ['30s hard timeout', 'force-release all', 'fires even if', 'event loop frozen'])

# ── Tier B: NETWORK (layers 4-6) ──
tier_band(s, 0.35, 3.2, 1.65, 'NETWORK', 'delivery guarantee', ORANGE)
ptt_layer(s, TIER_X, 3.2, cw3, 1.65, 4, 'ACK Retry', ORANGE,
          ['setPTT:false → 1s', 'await getPTT:false', 'retry 3× echo', 'ptt_manager.js'])
ptt_layer(s, TIER_X+cw3+CARD_GAP, 3.2, cw3, 1.65, 5, 'Backup Channel', PURPLE,
          ['wsAudioTX "s:"', 'independent socket', 'server forced-RX', 'separate from ctrl'])
ptt_layer(s, TIER_X+2*(cw3+CARD_GAP), 3.2, cw3, 1.65, 6, 'PING Health', RED,
          ['5s PING/PONG', 'timeout → dead', 's: + force reconnect', 'controls.js'])

# ── Tier C: FALLBACK (layers 7-8) ──
tier_band(s, 0.35, 5.05, 1.65, 'FALLBACK', 'always-RX', ACCENT)
cw2 = (TIER_W - CARD_GAP) / 2
ptt_layer(s, TIER_X, 5.05, cw2, 1.65, 7, 'Page Lifecycle', GREEN,
          ['visibilitychange · blur · pagehide → forcePTTReleaseIfActive()',
           'iOS suspends JS when backgrounded → touchend never delivered · tx_button.js'])
ptt_layer(s, TIER_X+cw2+CARD_GAP, 5.05, cw2, 1.65, 8, 'Server Authority', ACCENT,
          ['ws_ctrl() s: → radio.set_ptt(False) + dsp_proc.set_ptt(False)',
           'broadcast getPTT:false to ALL clients · always available · server.py'])

# guarantee banner
gb = box(s, 0.35, 6.92, 12.5, 0.42, f=RGBColor(GREEN[0]//5, GREEN[1]//5, GREEN[2]//5), br=GREEN)
p = gb.text_frame.paragraphs[0]
p.text = '🛡️  GUARANTEE:  every failure mode — lost packet · half-open socket · frozen tab · app backgrounded — has an independent path back to RX'
p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(10); p.font.color.rgb = GREEN; p.font.bold = True; p.font.name = 'Inter'
n(s, 14)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — Device Telemetry 0x1F00 (Ch 8 AD-011)
# ═══════════════════════════════════════════════════════════════════
s = slide()
hdr(s, 'Device Telemetry — 0x1F00 Packet Layout', 'Verified against real 40m drive sweep (expert_40m_drive.pcap) × external wattmeter + ExpertSDR3', idx=15, eyebrow='Deep Dive · AD-011')

box(s, 0.5, 1.4, 12.3, 1.0, f=CARD, br=ACCENT)
txt(s, 0.7, 1.45, 11.5, 0.25, '34-Byte Packet:  0xFF32 header (10B) + seq + flags + off16(u16) + off18(f32) + off22(f32) + off26(f32) + off30(f32)', sz=12, c=ACCENT, b=True)
txt(s, 0.7, 1.8, 11.5, 0.25, 'Sent continuously in RX and TX modes  ·  Device does NOT expose reverse power  →  cannot compute SWR  →  use external ATR-1000 for SWR', sz=10, c=DIM)

# 3 verified columns
fields = [
    ('✅  off30: f32', 'Forward Watts', 'PEP envelope · Direct float\nMonotonic: 28%→3W  71%→54W\n88%→83W  100%→101W\nNo cubic/linear fit needed', GREEN),
    ('✅  off16: u16/10', 'Supply Volts', '~136 (13.6V) at idle\nSags as power rises\n0W→13.6  50W→13.1  100W→12.9V\nCorrelation with W = -0.79', GREEN),
    ('✅  off18: f32', 'PA Temp °C', '~42°C · Barely moves\nMonitors PA thermal state\nduring extended TX sessions\n\nFrontend: getTXTelem:\nwatts,volts,temp,W_int', GREEN),
]
fx = 0.6
for offset, name, desc, clr in fields:
    box(s, fx, 2.65, 3.9, 2.2, f=CARD, br=clr)
    txt(s, fx+0.15, 2.73, 3.6, 0.25, offset, sz=14, c=clr, b=True)
    txt(s, fx+0.15, 3.05, 3.6, 0.25, name, sz=13, c=WHITE, b=True)
    txt(s, fx+0.15, 3.4, 3.6, 1.4, desc, sz=10, c=GRAY)
    fx += 4.1

# Wrong
box(s, 0.5, 5.1, 12.3, 1.7, f=CARD, br=RED)
txt(s, 0.7, 5.15, 6, 0.25, '⚠  Previously Wrong  (corrected 2026-06-25)', sz=14, c=ORANGE, b=True)
wrongs = [
    ('off14 u16 + cubic fit', '(pwr_raw-9)³×1.91e-5: non-monotonic noise, never tracked power'),
    ('off16 u16/100 as "SWR"', 'actually supply voltage — stable when NOT keyed → rules out SWR'),
    ('off22 f32 as "reverse power"', 'actually AVERAGE forward power (ratio to off30 ≈ SSB crest factor)'),
    ('off26 f32 as "SWR alt"', 'always exactly 1.0000 — device placeholder · 323+ packets identical'),
]
wx = 0.6
for title, desc in wrongs:
    box(s, wx, 5.55, 2.9, 1.1, f=CARD2, br=DIM)
    txt(s, wx+0.1, 5.6, 2.7, 0.5, title, sz=9, c=RED, b=True)
    txt(s, wx+0.1, 5.9, 2.7, 0.65, desc, sz=8, c=GRAY)
    wx += 3.1

txt(s, 0.8, 7.0, 11, 0.25, 'Frontend: getTXTelem:watts,volts,temp,W_int → UI displays W · VOLT (replaces former SWR) · TEMP', sz=9, c=DIM)
n(s, 15)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 16 — Version History + Thank You (Ch 14)
# ═══════════════════════════════════════════════════════════════════
s = slide()
txt(s, 0.8, 1.3, 11, 0.8, 'SunMRRC', sz=52, b=True)
txt(s, 0.8, 2.05, 11, 0.35, 'SunSDR2 DX  Mobile Remote Radio Control  ·  SDD V3.4', sz=18, c=ACCENT)
sep(s, 2.55)

txt(s, 0.8, 2.8, 11, 0.4, '📖  SDD Version Journey', sz=14, c=WHITE, b=True)
versions = [
    ('V1.0', '2026-03-15', 'Original MRRC TeamSD document set'),
    ('V2.0', '2026-05-02', 'MRRC V5.0 mobile UI and capability update'),
    ('V3.0', '2026-06-21', 'Re-baseline: FastAPI/Uvicorn, SunSDR2 DX direct UDP, explicit gaps'),
    ('V3.1', '2026-06-21', 'Waterfall rendering + S-meter smoothing'),
    ('V3.2', '2026-06-23', 'TX voice + DRIVE power + memory channel API documented'),
    ('V3.3', '2026-06-24', 'Dual-codec audio + SWR field correction + 11 ADs — V1.0 release'),
    ('V3.4', '2026-06-26', 'Telemetry correction · TX gain staging (AD-012) · 3 new SVG diagrams · PTT Ch 15 · 16 chapters'),
]
vy = 3.2
for ver, date, desc in versions:
    vc = GREEN if ver in ('V3.3', 'V3.4') else DIM
    tag(s, 0.9, vy, ver, vc)
    txt(s, 1.9, vy, 1.5, 0.22, date, sz=9, c=GRAY)
    txt(s, 3.5, vy, 9, 0.22, desc, sz=9.5, c=WHITE if ver in ('V3.3', 'V3.4') else GRAY)
    vy += 0.27

# Links + stats
txt(s, 0.8, 5.5, 11, 0.3, '🔗  Resources', sz=14, c=WHITE, b=True)
links = [
    ('🌐  Live Radio', 'https://radio.vlsc.net:8889'),
    ('📖  SDD Docs', 'https://www.vlsc.net/sunmrrc/sdd/'),
    ('🔧  GitHub', 'https://github.com/cheenle/sunsdr'),
    ('📥  PPTX Download', 'https://www.vlsc.net/sunmrrc/sdd/SunMRRC-Architecture-SDD-V3.4.pptx'),
    ('📐  SVG Diagrams', 'https://www.vlsc.net/sunmrrc/sdd/diagrams/'),
]
lx = 0.8
for label, url in links:
    txt(s, lx, 5.9, 2.4, 0.22, label, sz=10, c=WHITE, b=True)
    txt(s, lx, 6.15, 2.4, 0.2, url, sz=8, c=ACCENT)
    lx += 2.5

# Stats
box(s, 0.5, 6.5, 12.3, 0.5, f=CARD, br=DIM)
stats = [('16', 'SDD Chapters'), ('12', 'Architecture Decisions'), ('5', 'WebSocket Endpoints'),
         ('8', 'PTT Safety Layers'), ('5', 'SVG Diagrams'), ('4', 'TX EQ Presets'),
         ('195', 'TX IQ pkts/sec'), ('~20kbps', 'Opus Bandwidth')]
sx = 0.7
for nval, label in stats:
    txt(s, sx, 6.55, 1.2, 0.22, nval, sz=16, c=ACCENT, b=True, al=PP_ALIGN.CENTER)
    txt(s, sx, 6.75, 1.2, 0.2, label, sz=7, c=GRAY, al=PP_ALIGN.CENTER)
    sx += 1.5

txt(s, 0.8, 7.1, 11, 0.25, 'Document ID: SDD-SUNMRRC-2026-001  ·  2026-06-26  ·  Python 3.12 · FastAPI/Uvicorn · NumPy/SciPy · IBM TeamSD v2.3.2  ·  GPL-3.0', sz=8, c=DIM)
n(s, 16)

# ═══ Save ═════════════════════════════════════════════════════════
out = '/Users/cheenle/HAM/sunsdr/SDD/SunMRRC-Architecture-SDD-V3.4.pptx'
prs.save(out)
print(f'✅  {out}  ({len(prs.slides)} slides)')
